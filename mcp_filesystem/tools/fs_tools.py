"""文件系统工具 - fs_read, fs_write, fs_ops, fs_search.

整合前（7个工具）：
- read_file, read_multiple_files, write_file, create_directory,
  list_directory, move_file, get_file_info

整合后（4个工具）：
- fs_read: 读取文件（支持单文件/多文件，自动识别格式，包括 Excel）
- fs_write: 创建/覆盖文件（自动识别格式，包括 Excel）
- fs_ops: 文件系统操作（列目录、创建目录、移动、获取信息、删除、Excel元数据）
- fs_search: 搜索文件名或内容
"""

import json
import os
import re
from fnmatch import fnmatch
from functools import partial
from io import BytesIO
from pathlib import Path
from typing import Any, Dict, List, Literal, Optional, Tuple, Union

import anyio
from mcp.server.fastmcp.utilities.logging import get_logger

from ..operations import FileOperations, FileInfo
from ..security import PathValidator, sanitize_error_simple
from ..grep import GrepTools

logger = get_logger(__name__)

# 配置相关：读取 config.json 中的 fs 配置
CONFIG_PATH = Path(__file__).resolve().parents[2] / "config.json"
_config_cache: Optional[Dict[str, Any]] = None


def _load_config() -> Dict[str, Any]:
    """加载配置，带缓存。"""
    global _config_cache
    if _config_cache is not None:
        return _config_cache
    try:
        with open(CONFIG_PATH, "r", encoding="utf-8") as f:
            _config_cache = json.load(f)
            logger.info(f"Loaded fs config from {CONFIG_PATH}")
    except Exception as e:
        logger.warning(f"Failed to load fs config from {CONFIG_PATH}: {e}")
        _config_cache = {}
    return _config_cache


def _get_table_char_limit(default: int = 20000) -> int:
    """获取表格返回内容的最大字符数阈值（默认 20k）。"""
    cfg = _load_config().get("fs", {})
    raw_value = cfg.get("max_table_chars", default)
    try:
        value = int(raw_value)
        if value <= 0:
            raise ValueError("max_table_chars must be positive")
        return value
    except Exception:
        logger.warning(f"Invalid fs.max_table_chars '{raw_value}', fallback to {default}")
        return default


def _get_doc_char_limit(default: int = 20000) -> int:
    """获取文档（文本/Markdown）返回内容的最大字符数阈值（默认 20k）。"""
    cfg = _load_config().get("fs", {})
    raw_value = cfg.get("max_doc_chars", default)
    try:
        value = int(raw_value)
        if value <= 0:
            raise ValueError("max_doc_chars must be positive")
        return value
    except Exception:
        logger.warning(f"Invalid fs.max_doc_chars '{raw_value}', fallback to {default}")
        return default


def _parse_line_range(line_range: Optional[str]) -> (Optional[int], Optional[int]):
    """解析行区间字符串，如 '10-50', '10:', ':50', '20'."""
    if not line_range:
        return None, None
    s = str(line_range).strip()
    if not s:
        return None, None
    # 形式 "start-end" 或 "start:" 或 ":end"
    m = re.match(r"^\s*(\d+)?\s*[-:]\s*(\d+)?\s*$", s)
    if m:
        start = int(m.group(1)) if m.group(1) else None
        end = int(m.group(2)) if m.group(2) else None
        return start, end
    # 单一行号
    if s.isdigit():
        n = int(s)
        return n, n
    return None, None


def _normalize_mcp_param(value: Any, expected_type: type = None) -> Any:
    """规范化 MCP client 传递的参数.
    
    MCP client 可能将单个值错误包装成列表，此函数用于解包。
    
    Args:
        value: 原始参数值
        expected_type: 期望的类型（用于类型检查）
        
    Returns:
        规范化后的值
        
    Examples:
        _normalize_mcp_param(['hello']) -> 'hello'  # 单元素列表解包为字符串
        _normalize_mcp_param([['a', 'b']]) -> [['a', 'b']]  # 多元素列表保持原样
        _normalize_mcp_param('hello') -> 'hello'  # 非列表保持原样
    """
    # 如果是列表且只有一个元素
    if isinstance(value, list) and len(value) == 1:
        single_item = value[0]
        
        # 如果期望类型是字符串，且单元素是字符串，则解包
        if expected_type == str and isinstance(single_item, str):
            return single_item
        
        # 如果期望类型是字典，且单元素是字典，则解包
        if expected_type == dict and isinstance(single_item, dict):
            return single_item
        
        # 如果期望类型是列表，但单元素也是列表，需要判断：
        # - 如果单元素是列表的列表（2D数组），保持原样
        # - 如果单元素是普通列表，解包（可能是字符串列表被包装）
        if expected_type == list:
            if isinstance(single_item, list):
                # 检查是否是 2D 数组（列表的列表）
                if single_item and isinstance(single_item[0], list):
                    # 这是 2D 数组，保持原样
                    return value
                else:
                    # 这是普通列表，解包
                    return single_item
        
        # 其他情况：如果单元素类型匹配期望类型，解包
        if expected_type and isinstance(single_item, expected_type):
            return single_item
    
    return value


def _detect_format(path: Path) -> str:
    """根据文件扩展名自动检测格式.
    
    Returns:
        'text' | 'xlsx' | 'csv' | 'docx' | 'pptx' | 'pdf' | 'binary'
        
    Note: JSON、YAML 等格式统一识别为 'text'，使用统一的文本处理逻辑。
    """
    ext = path.suffix.lower()
    format_map = {
        '.xlsx': 'xlsx',
        '.xls': 'xlsx',
        '.csv': 'csv',
        '.docx': 'docx',
        '.doc': 'docx',
        '.pptx': 'pptx',
        '.ppt': 'pptx',
        '.pdf': 'pdf',
        # 常见文本格式
        '.txt': 'text',
        '.md': 'text',
        '.json': 'text',
        '.yaml': 'text',
        '.yml': 'text',
        '.py': 'text',
        '.js': 'text',
        '.ts': 'text',
        '.html': 'text',
        '.css': 'text',
        '.xml': 'text',
        '.sql': 'text',
        '.sh': 'text',
        '.bat': 'text',
        '.ini': 'text',
        '.conf': 'text',
        '.cfg': 'text',
        '.log': 'text',
        '.rst': 'text',
        '.tex': 'text',
    }
    return format_map.get(ext, 'text')  # 默认当作文本


def _is_binary_format(fmt: str) -> bool:
    """检查格式是否为需要回退到 .md 文件的二进制格式."""
    return fmt in ('docx', 'pptx', 'pdf', 'xlsx')


async def _try_fallback_to_md(
    abs_path: Path,
    virtual_path: str,
    validator: PathValidator,
) -> Optional[Tuple[Path, str]]:
    """尝试查找同名的 .md 文件作为回退.
    
    Args:
        abs_path: 原始文件的绝对路径
        virtual_path: 原始文件的虚拟路径
        validator: 路径验证器
        
    Returns:
        如果找到 .md 文件，返回 (md_abs_path, md_virtual_path) 的元组
        否则返回 None
    """
    # 构建同名的 .md 文件路径
    md_abs_path = abs_path.with_suffix('.md')
    
    # 检查 .md 文件是否存在且在允许的目录内
    if md_abs_path.exists() and md_abs_path.is_file():
        try:
            # 验证 .md 文件路径是否在允许的目录内
            if validator.is_path_allowed(md_abs_path):
                # 构建虚拟路径
                md_virtual_path = validator.real_to_virtual(md_abs_path)
                logger.info(
                    f"Found fallback .md file for '{virtual_path}': {md_virtual_path}"
                )
                return (md_abs_path, md_virtual_path)
        except Exception as e:
            logger.debug(f"Failed to validate fallback .md path: {e}")
    
    return None


def _iter_docx_blocks(doc):
    """按文档顺序遍历段落与表格。"""
    from docx.oxml.table import CT_Tbl
    from docx.oxml.text.paragraph import CT_P
    from docx.table import Table
    from docx.text.paragraph import Paragraph

    for child in doc.element.body:
        if isinstance(child, CT_P):
            yield ("paragraph", Paragraph(child, doc))
        elif isinstance(child, CT_Tbl):
            yield ("table", Table(child, doc))


def _docx_table_to_markdown(table):
    """将 docx 表格转换为 Markdown 表格。"""
    rows = []
    for row in table.rows:
        cells = []
        for cell in row.cells:
            cell_text = "\n".join(
                [p.text.strip() for p in cell.paragraphs if p.text and p.text.strip()]
            ).strip()
            cell_text = cell_text.replace("|", "\\|").replace("\n", " ")
            cells.append(cell_text)
        rows.append(cells)

    if not rows:
        return ""

    col_count = max(len(r) for r in rows)
    padded_rows = [r + [""] * (col_count - len(r)) for r in rows]

    header = padded_rows[0]
    separator = ["---"] * col_count
    body = padded_rows[1:]

    lines = [
        "| " + " | ".join(header) + " |",
        "| " + " | ".join(separator) + " |",
    ]
    for r in body:
        lines.append("| " + " | ".join(r) + " |")
    return "\n".join(lines)


def _docx_bytes_to_markdown(content_bytes: bytes) -> str:
    """将 docx 字节转换为 Markdown，保留表格。"""
    from docx import Document

    doc = Document(BytesIO(content_bytes))
    parts = []
    for kind, block in _iter_docx_blocks(doc):
        if kind == "table":
            table_md = _docx_table_to_markdown(block)
            if table_md.strip():
                parts.append(table_md)
        else:
            text = block.text.strip()
            if text:
                parts.append(text)
    return "\n\n".join(parts)


def _build_outline_from_md(text_md: str) -> List[Dict[str, Any]]:
    """从 markdown 文本构建大纲（标题及行号范围）。

    优先使用 # 开头的标题；若不存在，再尝试将整行粗体 (**xx**) 视为大纲节点，
    并推算节点的内容行区间。
    """
    lines = text_md.splitlines()
    headings: List[Dict[str, Any]] = []
    heading_positions: List[Dict[str, Any]] = []
    for idx, line in enumerate(lines, start=1):
        m = re.match(r"^(#{1,6})\s+(.*\S)", line.strip())
        if m:
            level = len(m.group(1))
            title = m.group(2).strip()
            heading_positions.append({"line": idx, "level": level, "title": title})
    total_lines = len(lines)
    if heading_positions:
        for i, h in enumerate(heading_positions):
            start = h["line"]
            end = heading_positions[i + 1]["line"] - 1 if i + 1 < len(heading_positions) else total_lines
            headings.append({
                "title": h["title"],
                "level": h["level"],
                "line_start": start,
                "line_end": end,
            })
        return headings

    # 无 # 标题时，尝试整行粗体作为大纲
    bold_outline: List[Dict[str, Any]] = []
    bold_lines: List[Dict[str, Any]] = []
    bold_pattern = re.compile(r"^\s*\*\*(.+?)\*\*\s*$")
    for idx, line in enumerate(lines, start=1):
        m = bold_pattern.match(line)
        if m:
            bold_lines.append({"line": idx, "title": m.group(1).strip()})

    if not bold_lines:
        return [{
            "title": "(document)",
            "level": 0,
            "line_start": 1,
            "line_end": total_lines,
        }]

    for i, h in enumerate(bold_lines):
        start_line = h["line"]
        end_line = bold_lines[i + 1]["line"] - 1 if i + 1 < len(bold_lines) else total_lines
        bold_outline.append({
            "title": h["title"],
            "level": 2,  # 粗体行视为二级
            "line_start": start_line,
            "line_end": end_line,
        })

    return bold_outline


async def _read_single_file(
    virtual_path: str,
    validator: PathValidator,
    operations: FileOperations,
    excel_ops: Any,
    *,
    sheet: Optional[str] = None,
    range: Optional[str] = None,
    max_rows: Optional[int] = None,
    encoding: str = "utf-8",
    line_start: Optional[int] = None,
    line_end: Optional[int] = None,
) -> Dict[str, Any]:
    """读取单个文件，自动识别格式.

    Args:
        virtual_path: Virtual path string (NOT a Path object)
    """
    logger.info(f"_read_single_file: Starting to read file: {virtual_path}")
    
    # Get real path for format detection
    try:
        abs_path, allowed = await validator.validate_path(virtual_path)
        if not allowed:
            error_msg = f"Path outside allowed directories: {virtual_path}"
            logger.error(f"_read_single_file: {error_msg}")
            raise ValueError(error_msg)
        logger.debug(f"_read_single_file: Validated path: {abs_path}")
    except Exception as e:
        logger.error(f"_read_single_file: Path validation failed for '{virtual_path}': {type(e).__name__}: {e}", exc_info=True)
        raise

    if not abs_path.exists():
        error_msg = f"File not found: {virtual_path} (real path: {abs_path})"
        logger.error(f"_read_single_file: {error_msg}")
        raise FileNotFoundError(error_msg)
    
    if not abs_path.is_file():
        error_msg = f"Path is not a file: {virtual_path} (real path: {abs_path})"
        logger.error(f"_read_single_file: {error_msg}")
        raise ValueError(error_msg)

    fmt = _detect_format(abs_path)
    logger.info(f"_read_single_file: Detected format: {fmt} for file: {virtual_path}")
    
    def _doc_response(text_md: str, fmt_label: str) -> Dict[str, Any]:
        """根据长度阈值或行范围返回内容/大纲。"""
        lines = text_md.splitlines()
        total_lines = len(lines)
        doc_limit = _get_doc_char_limit()

        # 按行读取
        if line_start is not None or line_end is not None:
            start = max(1, line_start or 1)
            end = min(total_lines, line_end if line_end is not None else total_lines)
            if end < start:
                start, end = end, start
            selected = lines[start - 1 : end]
            return {
                "format": fmt_label,
                "line_start": start,
                "line_end": end,
                "total_lines": total_lines,
                "content": "\n".join(selected),
            }

        # 未指定行范围，根据长度阈值自动控制
        if len(text_md) <= doc_limit:
            return {
                "format": fmt_label,
                "total_lines": total_lines,
                "content": text_md,
            }

        # 超出阈值，返回大纲（含行号范围）
        outline = _build_outline_from_md(text_md)
        return {
            "format": fmt_label,
            "total_lines": total_lines,
            "outline": outline,
            "truncated": True,
            "note": "Content omitted due to size; set line_range (e.g., '10-50') to fetch specific lines.",
        }

    if fmt == 'xlsx':
        # 使用 Excel 操作读取
        if excel_ops is None:
            error_msg = "Excel operations not available"
            logger.error(f"_read_single_file: {error_msg}")
            raise ValueError(error_msg)
        
        try:
            logger.info(f"_read_single_file: Reading Excel file: {virtual_path}, sheet={sheet}, range={range}, max_rows={max_rows}")
            result = await excel_ops.read_excel(
                virtual_path,  # Pass virtual path string
                sheet=sheet,
                range_str=range,
                max_rows=max_rows,
                output_format="markdown",
            )
            logger.info(f"_read_single_file: Successfully read Excel file: {virtual_path}")
            
            # 根据配置阈值控制表格内容长度，超过则仅返回前 20 行
            table_char_limit = _get_table_char_limit()
            if "path" in result:
                # 避免重复返回 path
                result.pop("path", None)
            try:
                serialized_len = len(json.dumps(result, ensure_ascii=False))
            except Exception as e:
                logger.warning(f"_read_single_file: Failed to measure Excel content length: {type(e).__name__}: {e}")
                serialized_len = 0
            
            if table_char_limit and serialized_len > table_char_limit:
                rows = result.get("rows") or []
                preview_rows = rows[:20] if isinstance(rows, list) else []
                result["rows"] = preview_rows
                result["returned_rows"] = len(preview_rows)
                result["truncated"] = True

                logger.info(
                    f"_read_single_file: Excel content exceeded {table_char_limit} chars "
                    f"(len={serialized_len}), returned first {len(preview_rows)} rows"
                )
            
            return {
                "path": virtual_path,
                "content": result,
            }
        except Exception as e:
            logger.error(f"_read_single_file: Failed to read Excel file '{virtual_path}': {type(e).__name__}: {e}", exc_info=True)
            raise

    elif fmt == 'csv':
        # 使用 Excel 操作读取 CSV
        if excel_ops is None:
            # 回退到文本读取
            logger.info(f"_read_single_file: Reading CSV as text (excel_ops not available): {virtual_path}")
            try:
                content = await operations.read_file(virtual_path, encoding=encoding)
                return {
                    "path": virtual_path,
                    "content": content,
                }
            except Exception as e:
                logger.error(f"_read_single_file: Failed to read CSV as text '{virtual_path}': {type(e).__name__}: {e}", exc_info=True)
                raise
        try:
            logger.info(f"_read_single_file: Reading CSV file: {virtual_path}, max_rows={max_rows}")
            result = await excel_ops.read_excel(
                virtual_path,  # Pass virtual path string
                max_rows=max_rows,
                output_format="markdown",
            )
            logger.info(f"_read_single_file: Successfully read CSV file: {virtual_path}")
            return {
                "path": virtual_path,
                "content": result,
            }
        except Exception as e:
            logger.error(f"_read_single_file: Failed to read CSV file '{virtual_path}': {type(e).__name__}: {e}", exc_info=True)
            raise

    elif fmt == 'docx':
        try:
            raw = await operations.read_file_binary(virtual_path)
            text_md = _docx_bytes_to_markdown(raw)
            return _doc_response(text_md, "docx")
        except ImportError:
            return {
                "path": virtual_path,
                "content": "",
                "error": "python-docx not installed"
            }
        except Exception as e:
            logger.error(f"_read_single_file: Failed to read DOCX '{virtual_path}': {type(e).__name__}: {e}", exc_info=True)
            raise
    
    elif fmt == 'pptx':
        try:
            from pptx import Presentation
        except ImportError:
            return {
                "path": virtual_path,
                "content": "",
                "error": "python-pptx not installed"
            }
        try:
            raw = await operations.read_file_binary(virtual_path)
            prs = Presentation(BytesIO(raw))
            slides = []
            for slide in prs.slides:
                texts = []
                for shape in slide.shapes:
                    if hasattr(shape, "text") and shape.text:
                        t = shape.text.strip()
                        if t:
                            texts.append(t)
                if texts:
                    slides.append("\n".join(texts))
            text_md = "\n\n---\n\n".join(slides)
            return _doc_response(text_md, "pptx")
        except Exception as e:
            logger.error(f"_read_single_file: Failed to read PPTX '{virtual_path}': {type(e).__name__}: {e}", exc_info=True)
            raise
    
    elif fmt == 'pdf':
        try:
            import pypdf
        except ImportError:
            return {
                "path": virtual_path,
                "content": "",
                "error": "pypdf not installed"
            }
        try:
            raw = await operations.read_file_binary(virtual_path)
            reader = pypdf.PdfReader(BytesIO(raw))
            parts = []
            for page in reader.pages:
                text = page.extract_text() or ""
                text = text.strip()
                if text:
                    parts.append(text)
            text_md = "\n\n---\n\n".join(parts)
            return _doc_response(text_md, "pdf")
        except Exception as e:
            logger.error(f"_read_single_file: Failed to read PDF '{virtual_path}': {type(e).__name__}: {e}", exc_info=True)
            raise
    
    else:
        # 默认文本读取
        content = await operations.read_file(virtual_path, encoding=encoding)
        # 文本/Markdown：支持行读取或按阈值返回大纲
        fmt_label = "md" if abs_path.suffix.lower() in {".md", ".markdown"} else "text"
        doc_resp = _doc_response(content, fmt_label)
        return {
            "path": virtual_path,
            "content": doc_resp,
        }


async def fs_read(
    path: Union[str, List[str]],
    validator: PathValidator,
    operations: FileOperations,
    excel_ops: Any = None,
    *,
    sheet: Optional[str] = None,
    range: Optional[str] = None,
    max_rows: Optional[int] = None,
    encoding: str = "utf-8",
    line_range: Optional[str] = None,
) -> Dict[str, Any]:
    """读取文件内容，自动识别格式.

    Args:
        path: 文件路径（支持单文件或文件路径列表）
        validator: 路径验证器
        operations: 文件操作实例
        excel_ops: Excel 操作实例（可选）
        sheet: xlsx 专用，工作表名称
        range: xlsx 专用，读取范围（如 A1:C10）
        max_rows: 最大行数限制
        encoding: 文本编码
        line_range: 文本/Markdown 按行读取范围，如"10:","20:50"

    Returns:
        单文件: {"path": str, "content": any}
        多文件: {"files": {path: result, ...}, "total": int, "success": int, "failed": int}

    Examples:
        # 读取单个文本文件
        fs_read("test.txt")

        # 读取 Excel 文件指定 sheet
        fs_read("data.xlsx", sheet="Sheet1", range="A1:D10")

        # 读取多个文件
        fs_read(["a.txt", "b.json", "c.xlsx"])
    """
    # 处理多文件情况
    if isinstance(path, list):
        results: Dict[str, Any] = {}
        errors: Dict[str, str] = {}
        
        line_start_val, line_end_val = _parse_line_range(line_range)

        for p in path:
            try:
                abs_path, allowed = await validator.validate_path(p)
                if not allowed:
                    errors[p] = f"Path outside allowed directories: {p}"
                    continue
                
                # 检测格式（在检查文件存在之前，以便判断是否需要回退）
                fmt = _detect_format(abs_path)
                
                # 如果文件不存在，且是二进制格式，尝试查找同名的 .md 文件
                if not abs_path.exists():
                    if _is_binary_format(fmt):
                        fallback_result = await _try_fallback_to_md(abs_path, p, validator)
                        if fallback_result:
                            md_abs_path, md_virtual_path = fallback_result
                            # 使用 .md 文件继续处理
                            abs_path = md_abs_path
                            p = md_virtual_path
                            logger.info(f"fs_read: Using fallback .md file for '{p}': {md_virtual_path}")
                        else:
                            errors[p] = f"File not found: {p}"
                            continue
                    else:
                        errors[p] = f"File not found: {p}"
                        continue

                # Pass virtual path string, not abs_path
                result = await _read_single_file(
                    p, validator, operations, excel_ops,
                    sheet=sheet, range=range, max_rows=max_rows, encoding=encoding,
                    line_start=line_start_val, line_end=line_end_val,
                )
                results[result["path"]] = result
            except Exception as e:
                errors[p] = sanitize_error_simple(e, p)

        return {
            "files": results,
            "errors": errors if errors else None,
            "total": len(path),
            "success": len(results),
            "failed": len(errors),
        }
    
    # 单文件处理
    line_start_val, line_end_val = _parse_line_range(line_range)
    logger.info(f"fs_read: Processing single file: {path}")
    try:
        abs_path, allowed = await validator.validate_path(path)
        if not allowed:
            error_msg = f"Path outside allowed directories: {path}"
            logger.error(f"fs_read: {error_msg}")
            raise ValueError(error_msg)
        
        # 检测格式（在检查文件存在之前，以便判断是否需要回退）
        fmt = _detect_format(abs_path)
        
        # 如果文件不存在，且是二进制格式，尝试查找同名的 .md 文件
        if not abs_path.exists():
            if _is_binary_format(fmt):
                fallback_result = await _try_fallback_to_md(abs_path, path, validator)
                if fallback_result:
                    md_abs_path, md_virtual_path = fallback_result
                    # 使用 .md 文件继续处理
                    abs_path = md_abs_path
                    path = md_virtual_path
                    logger.info(f"fs_read: Using fallback .md file: {md_virtual_path}")
                else:
                    error_msg = f"File not found: {path} (real path: {abs_path})"
                    logger.error(f"fs_read: {error_msg}")
                    raise FileNotFoundError(error_msg)
            else:
                error_msg = f"File not found: {path} (real path: {abs_path})"
                logger.error(f"fs_read: {error_msg}")
                raise FileNotFoundError(error_msg)

        # 支持目录读取：直接列出目录而不是按文件处理
        if abs_path.is_dir():
            logger.info(f"fs_read: Path is directory, listing: {path}")
            entries = await operations.list_directory(path)
            return {
                "path": path,
                "is_directory": True,
                "entries": entries,
            }

        # Pass virtual path string, not abs_path
        logger.info(f"fs_read: Calling _read_single_file for: {path}")
        result = await _read_single_file(
            path, validator, operations, excel_ops,
            sheet=sheet, range=range, max_rows=max_rows, encoding=encoding,
            line_start=line_start_val, line_end=line_end_val,
        )
        logger.info(f"fs_read: Successfully read file: {path}")
        return result
    except Exception as e:
        logger.error(f"fs_read: Error reading file '{path}': {type(e).__name__}: {e}", exc_info=True)
        raise


async def fs_write(
    path: str,
    content: Any,
    validator: PathValidator,
    operations: FileOperations,
    excel_ops: Any = None,
    *,
    overwrite: bool = False,
    append: bool = False,
    sheet: Optional[str] = None,
    encoding: str = "utf-8",
) -> Dict[str, Any]:
    """创建新文件或完全覆盖现有文件.
    
    自动识别格式并写入：
    - .xlsx/.xls: 写入 Excel 文件
    - .csv: 写入 CSV 文件
    - .json: 写入 JSON 文件（自动格式化）
    - 其他: 写入文本文件
    
    Args:
        path: 文件路径
        content: 文件内容
            - 对于 Excel/CSV: 应为 2D 数组 [[...], [...]]
            - 对于 JSON: 任意可序列化对象
            - 对于文本: 字符串
        validator: 路径验证器
        operations: 文件操作实例
        excel_ops: Excel 操作实例（可选）
        overwrite: 是否允许覆盖现有文件（默认 False）
        sheet: xlsx 专用，工作表名称
        encoding: 文本编码
        
    Returns:
        {"success": True, "path": str, "message": str, ...}
        
    Examples:
        # 写入文本文件
        fs_write("test.txt", "Hello, World!")
        
        # 写入 JSON 文件
        fs_write("config.json", {"key": "value"})
        
        # 写入 Excel 文件
        fs_write("data.xlsx", [["Name", "Age"], ["Alice", 30], ["Bob", 25]])
    """
    abs_path, allowed = await validator.validate_path(path)
    if not allowed:
        raise ValueError(f"Path outside allowed directories: {path}")
    
    # 检查是否存在（append 模式下允许追加，不检查 overwrite）
    if abs_path.exists() and not overwrite and not append:
        raise FileExistsError(
            f"File already exists: {path}. Use overwrite=True to replace it or append=True to append."
        )
    
    # 自动创建父目录
    if not abs_path.parent.exists():
        await anyio.to_thread.run_sync(
            lambda: abs_path.parent.mkdir(parents=True, exist_ok=True)
        )
    
    # 处理 content 参数：如果是列表且只有一个元素，提取该元素
    # 这可以处理 MCP client 可能将字符串包装成列表的情况
    if isinstance(content, list) and len(content) == 1:
        content = content[0]
    
    fmt = _detect_format(abs_path)
    virtual_path = validator.real_to_virtual(abs_path)
    
    # 规范化 content 参数（处理 MCP client 可能将值包装成列表的情况）
    if fmt in ('xlsx', 'csv'):
        # Excel/CSV 格式：期望是 2D 数组，不解包
        pass
    else:
        # 文本/JSON 格式：期望是字符串或对象，解包单元素列表
        content = _normalize_mcp_param(content, expected_type=str)
    
    if fmt == 'xlsx':
        if excel_ops is None:
            raise ValueError("Excel operations not available for .xlsx files")
        
        # content 应该是 2D 数组
        if not isinstance(content, list):
            raise ValueError("Content for Excel files must be a 2D array")
        
        # 追加模式：读取现有数据并合并
        if append and abs_path.exists():
            try:
                existing_data = await excel_ops.read_excel(
                    path,
                    sheet=sheet or "Sheet1",
                    output_format="markdown",  # 使用 markdown 格式，但主要用 rows 字段
                )
                # 合并数据（跳过表头，只追加数据行）
                if existing_data and "rows" in existing_data:
                    existing_rows = existing_data["rows"]
                    # 如果现有数据有表头，保留表头；否则直接追加所有行
                    if existing_rows:
                        # 合并：保留第一行作为表头（如果新数据也有表头，跳过新数据的表头）
                        merged_data = existing_rows.copy()
                        # 如果新数据的第一行与现有数据的表头相同，跳过新数据的表头
                        if content and len(existing_rows) > 0:
                            if len(content) > 0 and content[0] == existing_rows[0]:
                                merged_data.extend(content[1:])
                            else:
                                merged_data.extend(content)
                        else:
                            merged_data.extend(content)
                        content = merged_data
            except Exception as e:
                logger.warning(f"Failed to read existing Excel file for append: {e}, creating new file")
                # 如果读取失败，按新建文件处理
        
        result = await excel_ops.write_excel(
            path,  # Already a virtual path string
            data=content,
            sheet=sheet or "Sheet1",
            overwrite=True,  # 追加模式下总是覆盖（因为已经合并了数据）
        )
        # result 已包含 success, path, rows_written 等
        return result
    
    elif fmt == 'csv':
        # CSV 文件：如果 content 是字符串，直接按文本写入；如果是 2D 数组，格式化后写入
        if isinstance(content, str):
            # 字符串直接写入（纯文本 CSV）
            if append and abs_path.exists():
                # 追加模式：追加到文件末尾，确保有换行符
                existing = await anyio.to_thread.run_sync(
                    lambda: abs_path.read_text(encoding=encoding)
                )
                content = existing + ('\n' if not existing.endswith('\n') else '') + content
            await anyio.to_thread.run_sync(
                lambda: abs_path.write_text(content, encoding=encoding)
            )
            return {
                "success": True,
                "path": virtual_path,
                "size": len(content.encode(encoding)),
            }
        elif isinstance(content, list):
            # 2D 数组：格式化后写入
            lines = []
            for row in content:
                line = ','.join(
                    f'"{str(c)}"' if ',' in str(c) or '"' in str(c) else str(c)
                    for c in row
                )
                lines.append(line)
            csv_text = '\n'.join(lines)
            if append and abs_path.exists():
                # 追加模式：追加到文件末尾，确保有换行符
                existing = await anyio.to_thread.run_sync(
                    lambda: abs_path.read_text(encoding=encoding)
                )
                csv_text = existing + ('\n' if not existing.endswith('\n') else '') + csv_text
            await anyio.to_thread.run_sync(
                lambda: abs_path.write_text(csv_text, encoding=encoding)
            )
            return {
                "success": True,
                "path": virtual_path,
                "size": len(csv_text.encode(encoding)),
            }
        else:
            raise ValueError("CSV content must be a string or 2D array")
    
    elif fmt == 'json':
        if not isinstance(content, str):
            content = json.dumps(content, ensure_ascii=False, indent=2)
        # JSON 追加模式：按文本追加（因为 JSON 追加不太合理，但可以按文本处理）
        if append and abs_path.exists():
            existing = await anyio.to_thread.run_sync(
                lambda: abs_path.read_text(encoding=encoding)
            )
            content = existing + ('\n' if not existing.endswith('\n') else '') + content
        # 直接写入文件，路径已经验证过
        await anyio.to_thread.run_sync(
            lambda: abs_path.write_text(content, encoding=encoding)
        )
        return {
            "success": True,
            "path": virtual_path,
            "size": len(content.encode(encoding)),
        }

    else:
        # 文本文件
        if not isinstance(content, str):
            content = str(content)
        # 追加模式：追加到文件末尾，确保有换行符
        if append and abs_path.exists():
            existing = await anyio.to_thread.run_sync(
                lambda: abs_path.read_text(encoding=encoding)
            )
            content = existing + ('\n' if not existing.endswith('\n') else '') + content
        # 直接写入文件，路径已经验证过
        await anyio.to_thread.run_sync(
            lambda: abs_path.write_text(content, encoding=encoding)
        )
        return {
            "success": True,
            "path": virtual_path,
            "size": len(content.encode(encoding)),
        }


async def fs_ops(
    operation: Literal["list", "mkdir", "move", "info", "delete"],
    path: str,
    validator: PathValidator,
    operations: FileOperations,
    excel_ops: Any = None,
    *,
    destination: Optional[str] = None,
    recursive: bool = False,
) -> Dict[str, Any]:
    """文件系统操作（列目录、创建目录、移动文件、获取信息、删除）.

    Args:
        operation: 操作类型
            - 'list': 列出目录内容
            - 'mkdir': 创建目录
            - 'move': 移动文件/目录
            - 'info': 获取文件/目录信息
            - 'delete': 删除文件/目录
        path: 目标路径
        validator: 路径验证器
        operations: 文件操作实例
        destination: move 操作专用，目标路径
        recursive: 递归操作（用于 mkdir 和 delete）

    Returns:
        操作结果字典
        
    Examples:
        # 列出目录
        fs_ops("list", "/data")
        
        # 创建目录
        fs_ops("mkdir", "/data/new_folder")
        
        # 移动文件
        fs_ops("move", "/data/old.txt", destination="/data/new.txt")
        
        # 获取文件信息
        fs_ops("info", "/data/file.txt")
        
        # 删除文件
        fs_ops("delete", "/data/temp.txt")
    """
    abs_path, allowed = await validator.validate_path(path)
    if not allowed:
        raise ValueError(f"Path outside allowed directories: {path}")
    
    virtual_path = validator.real_to_virtual(abs_path)
    
    if operation == "list":
        if not abs_path.exists():
            raise FileNotFoundError(f"Directory not found: {path}")
        if not abs_path.is_dir():
            raise ValueError(f"Not a directory: {path}")

        async def _list_recursive(dir_path: str) -> Dict[str, Any]:
            """递归列出目录为树形结构。"""
            items = await operations.list_directory(dir_path)
            children = []
            for item in items:
                if item.get("is_directory"):
                    child_tree = await _list_recursive(item["path"])
                    item["children"] = child_tree.get("children", [])
                children.append(item)
            return {"children": children}

        if recursive:
            tree = await _list_recursive(path)
            return {
                "success": True,
                "action": "list",
                "path": virtual_path,
                "entries": tree.get("children", []),
                "recursive": True,
            }
        else:
            entries = await operations.list_directory(path)
            return {
                "success": True,
                "action": "list",
                "path": virtual_path,
                "entries": entries,
            }
    
    elif operation == "mkdir":
        # Pass virtual path string, not abs_path
        await operations.create_directory(path, parents=recursive)
        return {
            "success": True,
            "action": "mkdir",
            "path": virtual_path,
            "message": f"Directory created: {virtual_path}",
        }
    
    elif operation == "move":
        if destination is None:
            raise ValueError("destination is required for move operation")

        dest_abs, dest_allowed = await validator.validate_path(destination)
        if not dest_allowed:
            raise ValueError(f"Destination outside allowed directories: {destination}")

        if not abs_path.exists():
            raise FileNotFoundError(f"Source not found: {path}")

        # Pass virtual path strings, not abs_path/dest_abs
        await operations.move_file(path, destination)
        return {
            "success": True,
            "action": "move",
            "source": virtual_path,
            "destination": validator.real_to_virtual(dest_abs),
            "message": f"Moved {virtual_path} to {validator.real_to_virtual(dest_abs)}",
        }
    
    elif operation == "info":
        if not abs_path.exists():
            raise FileNotFoundError(f"Path not found: {path}")

        info = FileInfo(abs_path, virtual_path)
        result = {
            "success": True,
            **info.to_dict(),
        }

        # 如果是 Excel 文件，添加 Excel 元数据
        if excel_ops and abs_path.suffix.lower() in ['.xlsx', '.xls', '.csv']:
            try:
                excel_metadata = await excel_ops.get_workbook_metadata(path)
                result["excel_metadata"] = {
                    "sheets": excel_metadata.get("sheets", []),
                    "has_formulas": excel_metadata.get("has_formulas", False),
                    "total_sheets": excel_metadata.get("total_sheets", 0),
                }
            except Exception:
                # Excel 元数据获取失败，不影响基本信息返回
                pass

        return result
    
    elif operation == "delete":
        if not abs_path.exists():
            raise FileNotFoundError(f"Path not found: {path}")
        
        if abs_path.is_dir():
            if recursive:
                import shutil
                await anyio.to_thread.run_sync(partial(shutil.rmtree, abs_path))
            else:
                await anyio.to_thread.run_sync(abs_path.rmdir)
        else:
            await anyio.to_thread.run_sync(abs_path.unlink)
        
        return {
            "success": True,
            "action": "delete",
            "path": virtual_path,
            "message": f"Deleted: {virtual_path}",
        }
    
    else:
        raise ValueError(f"Unknown operation: {operation}")


async def fs_search(
    search_type: Literal["filename", "content"],
    pattern: str,
    validator: PathValidator,
    grep_tools: GrepTools,
    operations: FileOperations,
    *,
    path: Optional[str] = None,
    max_results: int = 100,
    case_sensitive: bool = False,
    is_regex: bool = False,
    context_lines: int = 2,
) -> Dict[str, Any]:
    """搜索文件名或文件内容.
    
    Args:
        search_type: 搜索类型
            - 'filename': 搜索文件名（支持 glob 模式）
            - 'content': 搜索文件内容（支持正则）
        pattern: 搜索模式
            - filename: glob 模式，如 "*.py", "test_*.txt"
            - content: 文本或正则表达式
        validator: 路径验证器
        grep_tools: Grep 工具实例
        operations: 文件操作实例
        path: 搜索路径（默认整个工作区）
        max_results: 最大结果数（默认 100）
        case_sensitive: 是否区分大小写（默认 False）
        is_regex: content 模式是否为正则（默认 False）
        context_lines: 返回匹配行前后的上下文行数（默认 2，仅对 content 搜索有效）
        
    Returns:
        {"matches": [...], "total": int, ...}
        
    Examples:
        # 搜索 Python 文件
        fs_search("filename", "*.py")
        
        # 搜索包含 "TODO" 的文件
        fs_search("content", "TODO")
        
        # 使用正则搜索
        fs_search("content", r"def\\s+\\w+\\(", is_regex=True)
    """
    # 确定搜索路径
    if path:
        abs_path, allowed = await validator.validate_path(path)
        if not allowed:
            raise ValueError(f"Path outside allowed directories: {path}")
    else:
        # 默认搜索路径：优先使用 virtual_root，对应虚拟 “/”
        if validator.virtual_root is not None:
            abs_path = validator.virtual_root
            path = "/"  # 用于后续 virtual_base 计算
        else:
            allowed_dirs = validator.get_allowed_dirs()
            if not allowed_dirs:
                raise ValueError("No allowed directories configured")
            abs_path = Path(allowed_dirs[0])
    
    if not abs_path.exists():
        raise FileNotFoundError(f"Search path not found: {path or str(abs_path)}")
    
    # 如果没有 virtual_root，直接返回空结果，避免暴露真实路径
    if validator.virtual_root is None:
        logger.warning("fs_search: virtual_root is None; return empty result to avoid leaking real paths")
        virtual_base = path if path else "/"
        if virtual_base and not virtual_base.startswith("/"):
            virtual_base = "/" + virtual_base
        return {
            "success": True,
            "search_type": search_type,
            "pattern": pattern,
            "path": virtual_base,
            "matches": [],
            "total": 0,
            "truncated": False,
        }
    
    # 确定虚拟路径基础路径（有 virtual_root）
    try:
        virtual_base = validator.real_to_virtual(abs_path, strict=False)
        # 如果返回了占位符，说明路径不在 virtual_root 下，使用 "/" 作为根
        if virtual_base == "/[path_not_available]":
            virtual_base = "/"
    except (ValueError, Exception) as e:
        logger.warning(f"Failed to convert search path to virtual: {e}, using '/'")
        virtual_base = "/"
    
    if search_type == "filename":
        # 文件名搜索（使用 glob）
        matches: List[Dict[str, Any]] = []
        count = 0
        
        # 简化的 glob 搜索
        async def search_dir(dir_path: Path, depth: int = 0) -> None:
            nonlocal count
            if count >= max_results or depth > 10:
                return
            
            try:
                entries = await anyio.to_thread.run_sync(list, dir_path.iterdir())
                for entry in entries:
                    if count >= max_results:
                        break
                    
                    try:
                        # Check if path is allowed using direct check (entry is a real Path)
                        if not validator.is_path_allowed(entry):
                            continue

                        name = entry.name
                        # 不区分大小写匹配
                        if not case_sensitive:
                            match = fnmatch(name.lower(), pattern.lower())
                        else:
                            match = fnmatch(name, pattern)

                        if match:
                            try:
                                virtual_entry = validator.real_to_virtual(entry, strict=False)
                                matches.append({
                                    "path": virtual_entry,
                                    "name": name,
                                    "is_dir": entry.is_dir(),
                                    "size": entry.stat().st_size if entry.is_file() else None,
                                })
                                count += 1
                            except (ValueError, Exception) as e:
                                # Skip files that can't be converted to virtual path
                                logger.debug(f"Skipping file {entry}: {e}")
                                continue
                        
                        # 递归搜索子目录
                        if entry.is_dir():
                            await search_dir(entry, depth + 1)
                    
                    except (PermissionError, FileNotFoundError, ValueError):
                        continue
            
            except (PermissionError, FileNotFoundError):
                pass
        
        await search_dir(abs_path)
        
        return {
            "success": True,
            "search_type": "filename",
            "pattern": pattern,
            "path": virtual_base,
            "matches": matches,
            "total": len(matches),
            "truncated": count >= max_results,
        }
    
    elif search_type == "content":
        # 内容搜索（使用 grep）
        # grep_files 需要虚拟路径字符串，使用 virtual_base
        logger.info(f"fs_search: content search, virtual_base={virtual_base}, pattern={pattern}, is_regex={is_regex}, case_sensitive={case_sensitive}, context_lines={context_lines}")
        try:
            result = await grep_tools.grep_files(
                path=virtual_base,
                pattern=pattern,
                is_regex=is_regex,
                case_sensitive=case_sensitive,
                max_results=max_results,
                recursive=True,
                context_lines=context_lines,  # 使用用户指定的上下文行数
            )
            logger.info(f"fs_search: grep_files returned {result.total_matches} matches from {result.files_searched} files")
            
            # Convert matches to dict with error handling
            matches_dict = []
            for i, m in enumerate(result.matches):
                try:
                    matches_dict.append(m.to_dict())
                except Exception as e:
                    logger.error(f"fs_search: Failed to convert match {i} to dict: {type(e).__name__}: {e}", exc_info=True)
            
            logger.info(f"fs_search: Successfully converted {len(matches_dict)} matches to dict")
            
            return {
                "success": True,
                "search_type": "content",
                "pattern": pattern,
                "path": virtual_base,
                "matches": matches_dict,
                "total": result.total_matches,
                "files_searched": result.files_searched,
                "truncated": result.total_matches >= max_results,
            }
        except Exception as e:
            logger.error(f"fs_search: grep_files failed for pattern '{pattern}': {type(e).__name__}: {e}", exc_info=True)
            raise
    
    else:
        raise ValueError(f"Unknown search type: {search_type}")

