import json
import os
import re
from io import BytesIO, StringIO
from pathlib import Path
from typing import Any, Dict, Optional, Tuple, List

import httpx
import pandas as pd
from fastmcp.utilities.logging import get_logger
from pptx import Presentation
from PyPDF2 import PdfReader  # type: ignore
from docx import Document

logger = get_logger(__name__)

# 尝试导入jieba用于中文分词，如果没有则使用简单分词
try:
    import jieba
    _HAS_JIEBA = True
except ImportError:
    _HAS_JIEBA = False
    logger.warning("jieba not found, using simple tokenization for Chinese text")

DEFAULT_BASE_URL = "http://localhost:3000"
ENV_BASE_URL = "KB_SEARCH_BASE_URL"
ENV_TOKEN = "KB_SEARCH_TOKEN"
CONFIG_PATH = Path(__file__).resolve().parents[2] / "config.json"

_config_cache: Optional[Dict[str, Any]] = None


def _load_config() -> Dict[str, Any]:
    global _config_cache
    if _config_cache is not None:
        return _config_cache
    try:
        with open(CONFIG_PATH, "r", encoding="utf-8") as f:
            _config_cache = json.load(f)
            logger.info(f"Loaded kb config from {CONFIG_PATH}")
    except Exception as e:
        logger.warning(f"Failed to load kb config from {CONFIG_PATH}: {e}")
        _config_cache = {}
    return _config_cache


def _get_kb_config(
    base_url: Optional[str] = None,
    token: Optional[str] = None,
) -> Tuple[str, Optional[str]]:
    """获取知识库配置（base_url 和 token）。
    
    优先级：
    - base_url: 参数 > 配置文件 kb.base_url > 环境变量 > 默认
    - token: 参数 > 配置文件 kb.token > 环境变量
    
    Returns:
        (base_url, token)
    """
    cfg = _load_config().get("kb", {})
    
    # 获取 base_url
    if base_url:
        final_base_url = base_url.rstrip("/")
        logger.debug(f"Using custom base_url: {final_base_url}")
    elif cfg.get("base_url"):
        final_base_url = str(cfg["base_url"]).rstrip("/")
        logger.info(f"Loaded base_url from config.json: {final_base_url}")
    elif os.environ.get(ENV_BASE_URL):
        final_base_url = os.environ.get(ENV_BASE_URL, "").rstrip("/")
        logger.info(f"Loaded base_url from env {ENV_BASE_URL}: {final_base_url}")
    else:
        final_base_url = DEFAULT_BASE_URL
        logger.info(f"Using default base_url: {final_base_url}")
    
    # 获取 token
    if token:
        final_token = token
        logger.debug("Using custom token")
    elif cfg.get("token"):
        final_token = str(cfg["token"])
        logger.info("Loaded token from config.json")
    elif os.environ.get(ENV_TOKEN):
        final_token = os.environ.get(ENV_TOKEN)
        logger.info(f"Loaded token from env {ENV_TOKEN}")
    else:
        final_token = None
        logger.warning("No token found in config or env")
    
    return final_base_url, final_token


def _get_max_read_chars(default: int = 20000) -> int:
    """从配置获取单次返回的最大字符数，默认 20k。"""
    cfg = _load_config().get("kb", {})
    raw_value = cfg.get("max_read_chars", default)
    try:
        value = int(raw_value)
        if value <= 0:
            raise ValueError("max_read_chars must be positive")
        return value
    except Exception:
        logger.warning(f"Invalid kb.max_read_chars '{raw_value}', fallback to {default}")
        return default


def _get_chunk_search_config() -> Dict[str, Any]:
    cfg = _load_config().get("kb", {})
    chunk_cfg = cfg.get("chunk_search", {})
    return {
        "limit": chunk_cfg.get("limit", 20000),
        "similarity": chunk_cfg.get("similarity", 0.4),
        "searchMode": chunk_cfg.get("searchMode", "embedding"),
    }


def _extract_keywords_from_pattern(pattern: str) -> List[str]:
    """从 glob pattern 中提取关键词。
    
    例如：'**/*RMS*终端*自动化*升级*' -> ['RMS', '终端', '自动化', '升级']
    """
    # 移除路径分隔符和通配符，提取中间的关键词
    # 匹配 *关键词* 或 关键词* 或 *关键词 的模式
    keywords = []
    
    # 先移除路径部分（**/ 或 */）
    text = re.sub(r'^\*\*?/', '', pattern)
    text = re.sub(r'/\*\*?$', '', text)
    
    # 提取被 * 包围或相邻的关键词
    # 匹配模式：*关键词* 或 关键词* 或 *关键词
    parts = re.split(r'\*+', text)
    for part in parts:
        part = part.strip()
        if part and len(part) >= 2:  # 至少2个字符才认为是关键词
            keywords.append(part)
    
    # 如果没有提取到关键词，返回空列表（保持原 pattern 搜索）
    return keywords


def _get_item_key(item: Dict[str, Any]) -> Optional[str]:
    """获取用于去重的唯一标识。"""
    # 优先使用 url，其次使用其他唯一字段
    return item.get("url") or item.get("path") or item.get("id") or str(item)


def _tokenize_text(text: str) -> List[str]:
    """对文本进行分词，支持中文和英文。
    
    如果安装了jieba，使用jieba进行中文分词；否则使用简单方法。
    """
    if not text:
        return []
    
    # 移除路径分隔符和通配符
    text = re.sub(r'[*/\\]', ' ', text)
    text = text.strip()
    
    if _HAS_JIEBA:
        # 使用jieba分词
        tokens = jieba.cut(text, cut_all=False)
        return [t.strip() for t in tokens if t.strip() and len(t.strip()) >= 1]
    else:
        # 简单分词：按中文字符、英文单词、数字分割
        # 匹配中文字符、英文单词、数字
        tokens = re.findall(r'[\u4e00-\u9fff]+|[a-zA-Z]+|\d+', text)
        return [t for t in tokens if len(t) >= 1]


def _score_item(item: Dict[str, Any], original_pattern: str, keywords: List[str]) -> float:
    """对搜索结果项进行评分。
    
    评分规则：
    1. 提取item的path（或url、name等字段）
    2. 对path进行分词
    3. 计算与原始pattern和关键词的匹配度
    
    Returns:
        评分（0-1之间，越高越相关）
    """
    # 获取用于评分的文本（优先使用path，其次url、name等）
    text = (
        item.get("path") or 
        item.get("url") or 
        item.get("name") or 
        item.get("title") or 
        ""
    )
    
    if not text:
        return 0.0
    
    # 对文本进行分词（不区分大小写）
    text_lower = text.lower()
    text_tokens = set(_tokenize_text(text_lower))
    
    # 对原始pattern进行分词
    pattern_lower = original_pattern.lower()
    pattern_tokens = set(_tokenize_text(pattern_lower))
    
    # 计算匹配度
    score = 0.0
    
    # 1. 完全匹配原始pattern（最高分）
    if original_pattern.lower() in text_lower:
        score += 0.5
    
    # 2. 关键词匹配度
    matched_keywords = 0
    for kw in keywords:
        kw_lower = kw.lower()
        # 完全匹配关键词
        if kw_lower in text_lower:
            matched_keywords += 1
            score += 0.3
        # 分词匹配
        elif any(kw_lower in token or token in kw_lower for token in text_tokens):
            matched_keywords += 1
            score += 0.15
    
    # 3. 分词交集匹配度
    common_tokens = pattern_tokens & text_tokens
    if pattern_tokens:
        token_match_ratio = len(common_tokens) / len(pattern_tokens)
        score += token_match_ratio * 0.2
    
    # 4. 关键词覆盖率
    if keywords:
        keyword_coverage = matched_keywords / len(keywords)
        score += keyword_coverage * 0.1
    
    # 归一化到0-1
    return min(1.0, score)


async def _kb_search_single(
    pattern: str,
    limit: int,
    base: str,
    auth_token: str,
    client: Optional[httpx.AsyncClient] = None,
) -> Dict[str, Any]:
    """执行单次知识库搜索。
    
    Args:
        pattern: glob pattern
        limit: 结果限制
        base: 基础 URL
        auth_token: 认证 token
        client: 可选的 HTTP 客户端（用于批量搜索时共享连接）
    """
    url = f"{base}/api/core/dataset/collection/pathMap"
    payload = {"globPattern": pattern, "limit": limit}
    headers = {
        "Authorization": f"Bearer {auth_token}",
        "Content-Type": "application/json",
    }

    use_external_client = client is not None
    created_client = None
    
    try:
        if not use_external_client:
            created_client = httpx.AsyncClient(timeout=30.0)
            client = created_client
        
        resp = await client.post(url, json=payload, headers=headers)

        status_code = resp.status_code
        try:
            data = resp.json()
        except Exception:
            data = {"raw": resp.text}

        if status_code >= 400:
            return {
                "success": False,
                "status_code": status_code,
                "error": data.get("message") if isinstance(data, dict) else resp.text,
            }

        if isinstance(data, dict):
            remote_code = data.get("code")
            if remote_code not in (None, 200):
                return {
                    "success": False,
                    "status_code": status_code,
                    "remote_code": remote_code,
                    "error": data.get("message") or "Remote returned non-200 code",
                }
            items = data.get("data", [])
        else:
            remote_code = None
            items = data

        return {
            "success": True,
            "status_code": status_code,
            "remote_code": remote_code,
            "items": items if isinstance(items, list) else [],
        }

    except httpx.HTTPError as e:
        logger.error(f"kb_search http error: {e}")
        return {"success": False, "error": f"HTTP error: {e}"}
    except Exception as e:
        logger.error(f"kb_search error: {e}", exc_info=True)
        return {"success": False, "error": str(e)}
    finally:
        # 只关闭我们自己创建的客户端
        if created_client:
            await created_client.aclose()


async def kb_search(
    pattern: str,
    limit: int = 10,
    *,
    base_url: Optional[str] = None,
    token: Optional[str] = None,
) -> Dict[str, Any]:
    """企业知识库 glob 搜索（不区分大小写），返回匹配文件列表。
    
    搜索策略：
    1. 先按原始 pattern 搜索
    2. 如果结果为空或少于 limit，提取关键词进行批量搜索
    3. 对搜索结果中的 path 进行分词，与原始输入匹配评分
    4. 按得分从高到低排序，返回 top limit 结果
    """
    try:
        limit_int = int(limit)
    except Exception:
        limit_int = 10
    limit_int = max(1, min(limit_int, 50))

    base, auth_token = _get_kb_config(base_url=base_url, token=token)

    if not auth_token:
        return {
            "success": False,
            "error": f"Missing auth token, set config.json kb.token or env {ENV_TOKEN}",
        }

    # 步骤1: 先按原始 pattern 搜索
    original_result = await _kb_search_single(pattern, limit_int, base, auth_token)
    
    original_items: List[Dict[str, Any]] = []
    if original_result.get("success") and original_result.get("items"):
        original_items = original_result["items"]
    
    # 步骤2: 如果结果不足，进行关键词拆分搜索
    all_items: List[Dict[str, Any]] = []
    seen_keys: set = set()
    
    # 添加原始搜索结果
    for item in original_items:
        key = _get_item_key(item)
        if key and key not in seen_keys:
            seen_keys.add(key)
            all_items.append(item)
    
    # 如果结果不足，提取关键词进行扩展搜索
    if len(all_items) < limit_int:
        keywords = _extract_keywords_from_pattern(pattern)
        
        if len(keywords) > 1:
            logger.info(
                f"Original search returned {len(all_items)} results (limit={limit_int}), "
                f"expanding with keywords: {keywords}"
            )
            
            # 对每个关键词分别搜索（使用共享的 HTTP 客户端以提高效率）
            search_patterns = [f"**/*{kw}*" for kw in keywords]
            
            async with httpx.AsyncClient(timeout=30.0) as client:
                for kw_pattern in search_patterns:
                    result = await _kb_search_single(kw_pattern, limit_int * 2, base, auth_token, client=client)
                    if result.get("success") and result.get("items"):
                        for item in result["items"]:
                            key = _get_item_key(item)
                            if key and key not in seen_keys:
                                seen_keys.add(key)
                                all_items.append(item)
                                # 如果已经收集足够的结果，可以提前停止
                                if len(all_items) >= limit_int * 3:  # 收集更多以便后续评分筛选
                                    break
                    if len(all_items) >= limit_int * 3:
                        break
        elif len(keywords) == 1:
            # 单个关键词，尝试更宽泛的搜索
            kw = keywords[0]
            logger.info(f"Single keyword detected: {kw}, trying broader search")
            result = await _kb_search_single(f"**/*{kw}*", limit_int * 2, base, auth_token)
            if result.get("success") and result.get("items"):
                for item in result["items"]:
                    key = _get_item_key(item)
                    if key and key not in seen_keys:
                        seen_keys.add(key)
                        all_items.append(item)
                        if len(all_items) >= limit_int * 3:
                            break
    
    # 步骤3: 对结果进行评分和排序
    if all_items:
        keywords = _extract_keywords_from_pattern(pattern)
        
        # 为每个item计算得分
        scored_items: List[Tuple[float, Dict[str, Any]]] = []
        for item in all_items:
            score = _score_item(item, pattern, keywords)
            scored_items.append((score, item))
        
        # 按得分从高到低排序
        scored_items.sort(key=lambda x: x[0], reverse=True)
        
        # 取前 limit_int 个结果
        top_items = [item for _, item in scored_items[:limit_int]]
        
        # 如果原始搜索结果已经足够，优先保留原始结果（即使得分可能不是最高）
        if len(original_items) > 0 and len(original_items) <= limit_int:
            # 将原始结果标记为高优先级（在得分相同的情况下）
            original_keys = {_get_item_key(item) for item in original_items}
            # 重新排序：原始结果优先，然后按得分
            scored_items_with_priority = []
            for score, item in scored_items:
                key = _get_item_key(item)
                priority = 1.0 if key in original_keys else 0.0
                # 原始结果得分加一个小偏移，确保优先
                adjusted_score = score + (priority * 0.01)
                scored_items_with_priority.append((adjusted_score, item))
            
            scored_items_with_priority.sort(key=lambda x: x[0], reverse=True)
            top_items = [item for _, item in scored_items_with_priority[:limit_int]]
        
        # 使用原始搜索的状态码
        status_code = original_result.get("status_code", 200)
        remote_code = original_result.get("remote_code", 200)
        
        meta = {
            "original_pattern": pattern,
            "original_results_count": len(original_items),
            "expanded_results_count": len(all_items),
        }
        if keywords:
            meta["extracted_keywords"] = keywords
        
        return {
            "success": True,
            "status_code": status_code,
            "remote_code": remote_code,
            "count": len(top_items),
            "items": top_items,
            "_meta": meta,
        }
    
    # 没有结果
    status_code = original_result.get("status_code", 200)
    remote_code = original_result.get("remote_code", 200)
    
    return {
        "success": True,
        "status_code": status_code,
        "remote_code": remote_code,
        "count": 0,
        "items": [],
        "_meta": {
            "original_pattern": pattern,
            "original_results_count": 0,
        },
    }


def _detect_ext(url: str, content_type: Optional[str]) -> str:
    from urllib.parse import urlparse

    ext = ""
    path = urlparse(url).path
    if "." in path:
        ext = path.rsplit(".", 1)[-1].lower()
    if not ext and content_type:
        if "pdf" in content_type:
            ext = "pdf"
        elif "msword" in content_type or "wordprocessingml" in content_type:
            ext = "docx"
        elif "presentationml" in content_type:
            ext = "pptx"
        elif "spreadsheetml" in content_type:
            ext = "xlsx"
        elif "csv" in content_type:
            ext = "csv"
        elif "markdown" in content_type:
            ext = "md"
        elif "html" in content_type:
            ext = "html"
        elif "json" in content_type:
            ext = "json"
        elif "text" in content_type:
            ext = "txt"
    return f".{ext}" if ext else ""


def _pdf_to_text(content: bytes) -> Tuple[str, bool]:
    reader = PdfReader(BytesIO(content))
    texts = []
    for page in reader.pages:
        try:
            texts.append(page.extract_text() or "")
        except Exception:
            texts.append("")
    text = "\n\n".join(t.strip() for t in texts if t is not None)
    has_text = bool(text.strip())
    return text.strip(), has_text


def _docx_to_text(content: bytes) -> str:
    doc = Document(BytesIO(content))
    parts = []
    for p in doc.paragraphs:
        if p.text.strip():
            parts.append(p.text)
    return "\n\n".join(parts)


def _pptx_to_text(content: bytes) -> str:
    prs = Presentation(BytesIO(content))
    parts = []
    for slide in prs.slides:
        for shape in slide.shapes:
            if hasattr(shape, "text") and shape.text:
                parts.append(shape.text)
    return "\n\n".join(parts)


def _csv_to_md(content: bytes) -> str:
    df = pd.read_csv(StringIO(content.decode("utf-8", errors="replace")))
    return df.to_markdown(index=False)


def _xlsx_to_md(content: bytes) -> str:
    # 取首个工作表，限制最大 200 行，避免过大输出
    df = pd.read_excel(BytesIO(content), sheet_name=0, nrows=200)
    return df.to_markdown(index=False)


async def kb_read_url(
    url: str,
    *,
    base_url: Optional[str] = None,
    token: Optional[str] = None,
    offset: int = 0,
    collection_ids: Optional[List[str]] = None,
    text: str = "",
) -> Dict[str, Any]:
    """读取知识库文件或按 collectionIds 做片段检索，返回 Markdown。
    
    - url: 直接下载并转为 Markdown，支持 offset 分页。
    - collection_ids + text: 调用后端向量检索，返回相关片段。
    单次返回长度由 config.json kb.max_read_chars 控制。

    TODO: 对于无文本的 PDF（扫描件），后续接入 OCR。
    """
    base, auth_token = _get_kb_config(base_url=base_url, token=token)

    full_url = url
    if not url.lower().startswith("http"):
        full_url = f"{base}{url if url.startswith('/') else '/' + url}"

    headers = {}
    if auth_token:
        headers["Authorization"] = f"Bearer {auth_token}"

    # 分支1：按 URL 读取完整文件
    if url:
        try:
            async with httpx.AsyncClient(timeout=60.0) as client:
                resp = await client.get(full_url, headers=headers)
            status_code = resp.status_code
            if status_code >= 400:
                return {
                    "success": False,
                    "status_code": status_code,
                    "error": f"HTTP {status_code}",
                }

            content_bytes = resp.content or b""
            content_type = resp.headers.get("content-type", "")
            size_bytes = len(content_bytes)

            # 尺寸保护：沿用 command.limits.max_file_size_mb，默认 100MB
            cfg = _load_config()
            max_mb = cfg.get("command", {}).get("limits", {}).get("max_file_size_mb", 100)
            max_bytes = max_mb * 1024 * 1024
            if size_bytes > max_bytes:
                return {
                    "success": False,
                    "error": f"File too large (> {max_mb}MB)",
                }

            ext = _detect_ext(full_url, content_type)
            text_md = ""
            todo_ocr = False

            if ext in [".md", ".txt", ".json", ".html", ".htm"]:
                text_md = content_bytes.decode("utf-8", errors="replace")
            elif ext == ".pdf":
                pdf_text, has_text = _pdf_to_text(content_bytes)
                text_md = pdf_text or ""
                if not has_text:
                    todo_ocr = True  # TODO: OCR for scanned PDFs
            elif ext == ".docx":
                text_md = _docx_to_text(content_bytes)
            elif ext == ".pptx":
                text_md = _pptx_to_text(content_bytes)
            elif ext == ".csv":
                text_md = _csv_to_md(content_bytes)
            elif ext in [".xlsx", ".xls"]:
                text_md = _xlsx_to_md(content_bytes)
            else:
                # 未识别类型，按文本尝试
                text_md = content_bytes.decode("utf-8", errors="replace")

            text_md = text_md or ""
            ext_label = ext or "unknown"

            # 分页裁剪：先全文转为 md，再按字符长度截取
            max_chars = _get_max_read_chars()
            try:
                start = int(offset)
            except Exception:
                start = 0
            start = max(0, start)

            total_chars = len(text_md)
            if start > total_chars:
                start = total_chars
            end = min(start + max_chars, total_chars)
            sliced_md = text_md[start:end]

            # 精简返回给智能体：仅保留必要信息
            result = {
                "success": True,
                "extension": ext_label,
                "content_md": sliced_md,
                "mode": "url",
            }

            needs_pagination = (start > 0) or (end < total_chars)
            if needs_pagination:
                page_index = max_chars and (start // max_chars + 1) or 1
                pagination = {
                    "page_start": start,
                    "page_end": end - 1 if end > 0 else 0,
                    "page_size": max_chars,
                    "returned_chars": len(sliced_md),
                    "total_chars": total_chars,
                    "has_more": end < total_chars,
                }
                if end < total_chars:
                    pagination["next_start"] = end
                    pagination["message"] = (
                        f"内容较长，已返回第 {page_index} 段。"
                        f" 继续读取可使用 offset={end} 再次调用 kb_read。"
                    )
                else:
                    pagination["message"] = "已返回最后一段内容。"
                result["pagination"] = pagination

            if todo_ocr:
                result["note"] = "PDF has no extractable text; OCR may be required."
            return result

        except httpx.HTTPError as e:
            logger.error(f"kb_read_url http error: {e}")
            return {"success": False, "error": f"HTTP error: {e}"}
        except Exception as e:
            logger.error(f"kb_read_url error: {e}", exc_info=True)
            return {"success": False, "error": str(e)}

    # 分支2：collectionIds + text 进行片段检索
    if collection_ids:
        if not text:
            return {"success": False, "error": "text is required when using collection_ids"}

        search_cfg = _get_chunk_search_config()
        payload_base = {
            "text": text,
            "searchMode": search_cfg["searchMode"],
            "limit": search_cfg["limit"],
            "similarity": search_cfg["similarity"],
        }
        all_items: List[Dict[str, Any]] = []
        try:
            async with httpx.AsyncClient(timeout=30.0) as client:
                for cid in collection_ids:
                    if not cid:
                        continue
                    payload = {**payload_base, "collectionId": cid}
                    resp = await client.post(
                        f"{base}/api/core/dataset/searchTest",
                        headers=headers,
                        json=payload,
                    )
                    if resp.status_code >= 400:
                        logger.warning(f"chunk search failed for {cid}: HTTP {resp.status_code}")
                        continue
                    try:
                        data = resp.json()
                    except Exception:
                        continue
                    if not isinstance(data, dict):
                        continue
                    if data.get("code") not in (None, 200):
                        continue
                    items = data.get("data", {}).get("list") or []
                    if isinstance(items, list):
                        all_items.extend(items)
        except httpx.HTTPError as e:
            logger.error(f"kb_read chunk search http error: {e}")
            return {"success": False, "error": f"HTTP error: {e}"}
        except Exception as e:
            logger.error(f"kb_read chunk search error: {e}", exc_info=True)
            return {"success": False, "error": str(e)}

        # 按检索顺序拼接片段，直到超出 chunk_search.limit（不分页，截断后丢弃后续）
        parts: List[str] = []
        used_chars = 0
        char_limit = search_cfg.get("limit", 20000)
        for idx, item in enumerate(all_items, start=1):
            score_val = None
            scores = item.get("score")
            if isinstance(scores, list) and scores:
                sc = scores[0]
                if isinstance(sc, dict):
                    score_val = sc.get("value")
            fragment = (
                f"[[片段:{idx}]] 来源:{item.get('sourceName') or item.get('collectionId') or ''} "
                f"(collectionId={item.get('collectionId', '')}, chunkIndex={item.get('chunkIndex', '')}, score={score_val})\n"
                f"{item.get('q', '')}"
            )
            if used_chars + len(fragment) > char_limit:
                break
            parts.append(fragment)
            used_chars += len(fragment) + 2  # 2 for the double newline join

        text_md = "\n\n".join(parts)

        result = {
            "success": True,
            "extension": "chunks",
            "content_md": text_md,
            "mode": "collection_chunks",
            "hits": len(all_items),
            "used_hits": len(parts),
            "char_limit": char_limit,
            "truncated": len(parts) < len(all_items),
        }

        return result

    # 无有效输入
    return {"success": False, "error": "Either url or collection_ids must be provided."}

