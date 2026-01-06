"""HTTP Routes for MCP Filesystem Server.

This module contains all HTTP API endpoints for:
- Admin API (workspace management, statistics)
- User API (file tree, downloads)
- Preview Deployment API
- Admin Web UI
"""

import json
import os
import shutil
import tempfile
import zipfile
from datetime import datetime
from io import BytesIO
from pathlib import Path
from typing import Any, Dict, Optional
from urllib.parse import unquote, urlparse
import mimetypes

import httpx
from fastmcp.utilities.logging import get_logger
from starlette.applications import Starlette
from starlette.requests import Request
from starlette.responses import FileResponse, JSONResponse, Response, StreamingResponse
from starlette.routing import Route

from .context import get_workspace_name
from .command.preview import PreviewManager
from .operations import format_timestamp

logger = get_logger(__name__)


# ========== Config and Auth ==========

def load_config() -> Dict[str, Any]:
    """Load configuration from config.json file."""
    config_path = Path(__file__).parent.parent / "config.json"
    if config_path.exists():
        try:
            with open(config_path, "r", encoding="utf-8") as f:
                return json.load(f)
        except Exception as e:
            logger.warning(f"Failed to load config: {e}")
    return {}


def get_admin_token() -> Optional[str]:
    """Get the admin token from config."""
    config = load_config()
    return config.get("admin_token")


async def healthcheck(request: Request):
    """Simple health endpoint for probes."""
    return JSONResponse({
        "success": True,
        "status": "ok",
        "service": "mcp-filesystem",
    })


async def workspace_viewer(request: Request):
    """Serve the workspace viewer HTML page."""
    static_dir = Path(__file__).parent.parent / "static"
    viewer_path = static_dir / "workspace_viewer.html"
    
    if not viewer_path.exists():
        return JSONResponse({
            "success": False,
            "error": "Workspace viewer not found"
        }, status_code=404)
    
    return FileResponse(viewer_path, media_type="text/html")


async def portal_homepage(request: Request):
    """Serve the project homepage."""
    portal_dir = Path(__file__).parent.parent / "portal"
    index_path = portal_dir / "index.html"
    
    if not index_path.exists():
        return JSONResponse({
            "success": False,
            "error": "Homepage not found"
        }, status_code=404)
    
    return FileResponse(index_path, media_type="text/html")


async def verify_admin_token(request: Request) -> bool:
    """Verify Bearer token authentication for admin endpoints."""
    expected_token = get_admin_token()
    if not expected_token:
        return False

    auth_header = request.headers.get("authorization", "")
    if auth_header.startswith("Bearer "):
        token = auth_header[7:]
        return token == expected_token
    return False


def get_user_data_dir() -> Path:
    """Get the user data directory path."""
    default_user_data_dir = Path(__file__).parent.parent / "user_data"
    return Path(os.environ.get("MCP_WORKSPACES_DIR", str(default_user_data_dir)))


def get_web_config() -> Dict[str, Any]:
    """Get web UI configuration."""
    cfg = load_config()
    return cfg.get("admin_web", {})


def human_size(size: int) -> str:
    """Convert bytes to human readable format."""
    units = ["B", "KB", "MB", "GB", "TB", "PB"]
    size_float = float(size)
    for unit in units:
        if size_float < 1024 or unit == units[-1]:
            return f"{size_float:.2f} {unit}"
        size_float /= 1024
    return f"{size_float:.2f} PB"


def get_folder_stats(folder_path: Path) -> Dict[str, Any]:
    """Get statistics for a folder."""
    try:
        stat = folder_path.stat()
        file_count = 0
        dir_count = 0
        total_size = 0

        for item in folder_path.rglob("*"):
            if item.is_file():
                file_count += 1
                total_size += item.stat().st_size
            elif item.is_dir():
                dir_count += 1

        return {
            "path": str(folder_path),
            "name": folder_path.name,
            "created_time": format_timestamp(stat.st_ctime),
            "modified_time": format_timestamp(stat.st_mtime),
            "accessed_time": format_timestamp(stat.st_atime),
            "file_count": file_count,
            "dir_count": dir_count,
            "total_size": total_size,
        }
    except Exception as e:
        return {"error": str(e)}


# ========== Admin API Endpoints ==========

async def admin_stats(request: Request):
    """Get overall statistics about all workspaces."""
    if not await verify_admin_token(request):
        return JSONResponse({"success": False, "error": "Unauthorized"}, status_code=401)

    user_data_dir = get_user_data_dir()
    if not user_data_dir.exists():
        return JSONResponse({
            "success": True,
            "user_data_dir": str(user_data_dir),
            "total_workspaces": 0,
            "unique_users": 0,
            "total_size_bytes": 0,
            "total_size_human": "0 B",
        })

    workspaces = [d for d in user_data_dir.iterdir() if d.is_dir()]
    unique_users = set()
    total_size = 0

    for ws in workspaces:
        name = ws.name
        if "_" in name:
            user_part = name.rsplit("_", 1)[0]
            unique_users.add(user_part)
        else:
            unique_users.add(name)

        for f in ws.rglob("*"):
            if f.is_file():
                total_size += f.stat().st_size

    return JSONResponse({
        "success": True,
        "user_data_dir": str(user_data_dir),
        "total_workspaces": len(workspaces),
        "unique_users": len(unique_users),
        "total_size_bytes": total_size,
        "total_size_human": human_size(total_size),
    })


async def admin_list_workspaces(request: Request):
    """List all workspaces with optional user_id filter."""
    if not await verify_admin_token(request):
        return JSONResponse({"success": False, "error": "Unauthorized"}, status_code=401)

    user_id_filter = request.query_params.get("user_id")
    user_data_dir = get_user_data_dir()

    if not user_data_dir.exists():
        return JSONResponse({"success": True, "workspaces": []})

    workspaces = []
    for ws in user_data_dir.iterdir():
        if not ws.is_dir():
            continue
        if user_id_filter and not ws.name.startswith(user_id_filter):
            continue
        workspaces.append(get_folder_stats(ws))

    return JSONResponse({"success": True, "workspaces": workspaces})


async def admin_workspace_info(request: Request):
    """Get information about a specific workspace."""
    if not await verify_admin_token(request):
        return JSONResponse({"success": False, "error": "Unauthorized"}, status_code=401)

    workspace_id = request.path_params.get("workspace_id")
    user_data_dir = get_user_data_dir()
    ws_path = user_data_dir / workspace_id

    if not ws_path.exists() or not ws_path.is_dir():
        return JSONResponse({"success": False, "error": "Workspace not found"}, status_code=404)

    return JSONResponse({"success": True, "workspace": get_folder_stats(ws_path)})


async def admin_workspace_tree(request: Request):
    """Get file tree for a workspace."""
    if not await verify_admin_token(request):
        return JSONResponse({"success": False, "error": "Unauthorized"}, status_code=401)

    workspace_id = request.path_params.get("workspace_id")
    max_depth = int(request.query_params.get("max_depth", "5"))
    user_data_dir = get_user_data_dir()
    ws_path = user_data_dir / workspace_id

    if not ws_path.exists() or not ws_path.is_dir():
        return JSONResponse({"success": False, "error": "Workspace not found"}, status_code=404)

    def build_tree(path: Path, depth: int = 0):
        if depth > max_depth:
            return {"name": path.name, "type": "directory", "truncated": True}

        if path.is_file():
            return {
                "name": path.name,
                "type": "file",
                "size": path.stat().st_size,
            }

        children = []
        try:
            for item in sorted(path.iterdir(), key=lambda x: (x.is_file(), x.name)):
                children.append(build_tree(item, depth + 1))
        except PermissionError:
            pass

        return {
            "name": path.name,
            "type": "directory",
            "children": children,
        }

    return JSONResponse({"success": True, "tree": build_tree(ws_path)})


async def admin_delete_workspace(request: Request):
    """Delete a workspace (requires confirmation)."""
    if not await verify_admin_token(request):
        return JSONResponse({"success": False, "error": "Unauthorized"}, status_code=401)

    workspace_id = request.path_params.get("workspace_id")
    confirm = request.query_params.get("confirm")

    if confirm != "yes":
        return JSONResponse({
            "success": False,
            "error": "Add ?confirm=yes to confirm deletion"
        }, status_code=400)

    user_data_dir = get_user_data_dir()
    ws_path = user_data_dir / workspace_id

    if not ws_path.exists() or not ws_path.is_dir():
        return JSONResponse({"success": False, "error": "Workspace not found"}, status_code=404)

    shutil.rmtree(ws_path)
    return JSONResponse({"success": True, "message": f"Workspace {workspace_id} deleted"})


async def admin_clear_cache(request: Request):
    """Clear the components cache (useful for debugging)."""
    if not await verify_admin_token(request):
        return JSONResponse({"success": False, "error": "Unauthorized"}, status_code=401)

    # Get components cache
    _components_cache = get_components_cache()
    cache_size = len(_components_cache)
    _components_cache.clear()

    return JSONResponse({
        "success": True,
        "message": f"Cleared {cache_size} cached workspace components",
        "cache_size_before": cache_size
    })


async def admin_preview_stats(request: Request):
    """Get all preview statistics (access count, creation time, last access time).
    
    Query parameters:
        sort: Sort field (last_accessed_at, created_at, access_count, url, workspace_name)
              Default: last_accessed_at
        order: Sort order (asc, desc). Default: desc
    """
    if not await verify_admin_token(request):
        return JSONResponse({"success": False, "error": "Unauthorized"}, status_code=401)

    try:
        from .command.preview import get_global_usage_tracker
        
        usage_tracker = get_global_usage_tracker()
        all_previews = usage_tracker.get_all_previews()
        
        # Get sort parameters
        sort_field = request.query_params.get("sort", "last_accessed_at")
        sort_order = request.query_params.get("order", "desc").lower()
        
        # Valid sort fields
        valid_sort_fields = {
            "last_accessed_at", "created_at", "access_count", 
            "url", "workspace_name", "subdomain"
        }
        if sort_field not in valid_sort_fields:
            return JSONResponse({
                "success": False,
                "error": f"Invalid sort field. Valid fields: {', '.join(valid_sort_fields)}"
            }, status_code=400)
        
        if sort_order not in ("asc", "desc"):
            return JSONResponse({
                "success": False,
                "error": "Invalid sort order. Use 'asc' or 'desc'"
            }, status_code=400)
        
        # Convert to list of dicts with formatted timestamps
        previews_data = []
        for preview in all_previews:
            preview_dict = preview.to_dict()
            
            # Format timestamps to readable format
            if preview_dict.get("created_at"):
                preview_dict["created_at_formatted"] = format_timestamp(preview_dict["created_at"])
            else:
                preview_dict["created_at_formatted"] = None
            
            if preview_dict.get("last_accessed_at"):
                preview_dict["last_accessed_at_formatted"] = format_timestamp(preview_dict["last_accessed_at"])
            else:
                preview_dict["last_accessed_at_formatted"] = None
            
            previews_data.append(preview_dict)
        
        # Sort the data
        reverse = (sort_order == "desc")
        
        def get_sort_key(item):
            value = item.get(sort_field)
            # Handle None values - put them at the end
            if value is None:
                # For descending order, None should be last (smallest)
                # For ascending order, None should be last (largest)
                if sort_field in ("access_count", "created_at", "last_accessed_at"):
                    return (0, 0)  # Numeric None
                return (0, "")  # String None
            
            # For numeric fields
            if sort_field in ("access_count", "created_at", "last_accessed_at"):
                return (1, value)
            # For string fields
            return (1, str(value).lower())
        
        previews_data.sort(key=get_sort_key, reverse=reverse)
        
        return JSONResponse({
            "success": True,
            "total": len(previews_data),
            "sort": {
                "field": sort_field,
                "order": sort_order
            },
            "previews": previews_data
        })
        
    except Exception as e:
        logger.error(f"Error getting preview stats: {e}", exc_info=True)
        return JSONResponse({
            "success": False,
            "error": str(e)
        }, status_code=500)


# ========== User API Endpoints ==========

async def user_workspace_tree(request: Request):
    """Get file tree for user's workspace (scoped by user_id + chat_id)."""
    MAX_ITEMS_PER_DIR = 20

    user_id = request.query_params.get("user_id") or request.headers.get("x-user-id")
    chat_id = request.query_params.get("chat_id") or request.headers.get("x-chat-id")
    max_depth_param = request.query_params.get("max_depth")

    try:
        max_depth = int(max_depth_param) if max_depth_param is not None else 5
        max_depth = max(0, max_depth)
    except ValueError:
        return JSONResponse({
            "success": False,
            "error": "Invalid max_depth value",
        }, status_code=400)

    if not user_id or not chat_id:
        return JSONResponse({
            "success": False,
            "error": "Missing required parameters: user_id and chat_id are required"
        }, status_code=400)

    workspace_name = get_workspace_name(user_id, chat_id)
    if not workspace_name:
        return JSONResponse({
            "success": False,
            "error": "Invalid user_id and chat_id combination"
        }, status_code=400)

    user_data_dir = get_user_data_dir().resolve()
    ws_path = (user_data_dir / workspace_name).resolve()

    # Security: ensure path is within user data directory
    try:
        ws_path.relative_to(user_data_dir)
    except ValueError:
        return JSONResponse({
            "success": False,
            "error": "Invalid workspace path"
        }, status_code=400)

    def build_tree(path: Path, depth: int = 0, stat_result=None):
        """Build tree with per-directory item caps ordered by modified time."""
        try:
            stat_result = stat_result or path.stat()
        except (FileNotFoundError, PermissionError):
            return None

        modified_time = format_timestamp(stat_result.st_mtime)

        if path.is_file() or path.is_symlink():
            return {
                "name": path.name,
                "type": "file",
                "size": stat_result.st_size,
                "modified_time": modified_time,
            }

        if depth > max_depth:
            return {
                "name": path.name,
                "type": "directory",
                "children": [],
                "modified_time": modified_time,
                "truncated": True,
            }

        entries_with_stats = []
        try:
            for child in path.iterdir():
                try:
                    child_stat = child.stat()
                    entries_with_stats.append((child, child_stat))
                except (FileNotFoundError, PermissionError):
                    continue
        except PermissionError:
            pass

        # Sort by last modified time desc, then name for stability
        entries_with_stats.sort(
            key=lambda pair: (-pair[1].st_mtime, pair[0].name)
        )

        children = []
        for child_path, child_stat in entries_with_stats[:MAX_ITEMS_PER_DIR]:
            child_node = build_tree(child_path, depth + 1, child_stat)
            if child_node:
                children.append(child_node)

        return {
            "name": path.name,
            "type": "directory",
            "children": children,
            "modified_time": modified_time,
        }

    # If workspace doesn't exist, return empty directory structure
    if not ws_path.exists() or not ws_path.is_dir():
        import time
        tree = {
            "name": workspace_name,
            "type": "directory",
            "children": [],
            "modified_time": format_timestamp(time.time()),
        }
    else:
        tree = build_tree(ws_path)

    return JSONResponse({
        "success": True,
        "tree": tree,
        "max_depth": max_depth,
        "per_directory_limit": MAX_ITEMS_PER_DIR,
    })


async def api_workspace_file_download(request: Request):
    """Download a single file from workspace."""
    try:
        user_id = request.query_params.get("user_id") or request.headers.get("x-user-id")
        chat_id = request.query_params.get("chat_id") or request.headers.get("x-chat-id")
        file_path = request.query_params.get("file_path")

        if not user_id or not chat_id:
            return JSONResponse({
                "success": False,
                "error": "Missing required parameters: user_id and chat_id are required"
            }, status_code=400)

        if not file_path:
            return JSONResponse({
                "success": False,
                "error": "Missing required parameter: file_path"
            }, status_code=400)

        workspace_name = get_workspace_name(user_id, chat_id)
        if not workspace_name:
            return JSONResponse({
                "success": False,
                "error": "Invalid user_id and chat_id combination"
            }, status_code=400)

        user_data_dir = get_user_data_dir().resolve()
        ws_path = (user_data_dir / workspace_name).resolve()

        # Security check
        try:
            ws_path.relative_to(user_data_dir)
        except ValueError:
            return JSONResponse({
                "success": False,
                "error": "Invalid workspace path"
            }, status_code=400)

        if not ws_path.exists() or not ws_path.is_dir():
            return JSONResponse({
                "success": False,
                "error": "Workspace not found"
            }, status_code=404)

        # Build target file path
        normalized_path = file_path.lstrip("/").replace("\\", "/")
        target_file = (ws_path / normalized_path).resolve()

        # Security check: ensure file is within workspace
        try:
            target_file.relative_to(ws_path)
        except ValueError:
            return JSONResponse({
                "success": False,
                "error": "Access denied: file path outside workspace"
            }, status_code=403)

        if not target_file.exists():
            return JSONResponse({
                "success": False,
                "error": "File not found"
            }, status_code=404)

        if not target_file.is_file():
            return JSONResponse({
                "success": False,
                "error": "Path is not a file"
            }, status_code=400)

        # Check file size limit
        config = load_config()
        max_file_size_mb = config.get("command", {}).get("limits", {}).get("max_file_size_mb", 100)
        max_file_size_bytes = max_file_size_mb * 1024 * 1024

        file_size = target_file.stat().st_size
        if file_size > max_file_size_bytes:
            return JSONResponse({
                "success": False,
                "error": f"File too large. Maximum size: {max_file_size_mb}MB"
            }, status_code=413)

        # Determine content type
        content_type = "application/octet-stream"
        if target_file.suffix:
            suffix_lower = target_file.suffix.lower()
            content_type_map = {
                ".txt": "text/plain",
                ".html": "text/html",
                ".css": "text/css",
                ".js": "text/javascript",
                ".json": "application/json",
                ".xml": "application/xml",
                ".pdf": "application/pdf",
                ".zip": "application/zip",
                ".jpg": "image/jpeg",
                ".jpeg": "image/jpeg",
                ".png": "image/png",
                ".gif": "image/gif",
            }
            content_type = content_type_map.get(suffix_lower, "application/octet-stream")

        return FileResponse(
            path=str(target_file),
            filename=target_file.name,
            media_type=content_type,
        )

    except Exception as e:
        logger.error(f"Error in workspace file download API: {e}", exc_info=True)
        return JSONResponse({
            "success": False,
            "error": f"Internal server error: {str(e)}"
        }, status_code=500)


TEXT_EXTENSIONS = {
    ".txt", ".md", ".markdown", ".json", ".yaml", ".yml", ".html", ".htm",
    ".css", ".js", ".ts", ".tsx", ".py", ".ini", ".toml", ".env", ".xml",
}


def _iter_docx_blocks(doc):
    """按文档顺序遍历段落与表格。"""
    from docx.oxml.table import CT_Tbl
    from docx.oxml.text.paragraph import CT_P
    from docx.table import Table
    from docx.text.paragraph import Paragraph

    for child in doc.element.body:
        if isinstance(child, CT_P):
            yield ("paragraph", Paragraph(child, doc))
        elif isinstance(child, CT_Tbl):
            yield ("table", Table(child, doc))


def _docx_table_to_markdown(table):
    """将 docx 表格转换为 Markdown 表格。"""
    rows = []
    for row in table.rows:
        cells = []
        for cell in row.cells:
            cell_text = "\n".join(
                [p.text.strip() for p in cell.paragraphs if p.text and p.text.strip()]
            ).strip()
            # 简单转义 |，保持单行
            cell_text = cell_text.replace("|", "\\|").replace("\n", " ")
            cells.append(cell_text)
        rows.append(cells)

    if not rows:
        return ""

    col_count = max(len(r) for r in rows)
    padded_rows = [r + [""] * (col_count - len(r)) for r in rows]

    header = padded_rows[0]
    separator = ["---"] * col_count
    body = padded_rows[1:]

    lines = [
        "| " + " | ".join(header) + " |",
        "| " + " | ".join(separator) + " |",
    ]
    for r in body:
        lines.append("| " + " | ".join(r) + " |")
    return "\n".join(lines)


def _docx_bytes_to_markdown(content_bytes: bytes) -> str:
    """将 docx 字节转换为 Markdown，保留表格。"""
    from docx import Document

    doc = Document(BytesIO(content_bytes))
    parts = []
    for kind, block in _iter_docx_blocks(doc):
        if kind == "table":
            table_md = _docx_table_to_markdown(block)
            if table_md.strip():
                parts.append(table_md)
        else:
            text = block.text.strip()
            if text:
                parts.append(text)
    return "\n\n".join(parts)


def _bytes_to_markdown(content_bytes: bytes, suffix: str, max_rows: int = 200) -> str:
    """将文件二进制内容转换为 Markdown 文本。
    
    Args:
        content_bytes: 文件的二进制内容
        suffix: 文件后缀名
        max_rows: 最大读取行数，-1 表示读取全部，默认 200 行（向后兼容）
    """
    ext = (suffix or "").lower()
    
    # 处理 max_rows 参数：只有 -1 表示读取全部
    nrows = None if max_rows == -1 else max_rows

    if ext in TEXT_EXTENSIONS or not ext:
        return content_bytes.decode("utf-8", errors="replace")

    if ext in {".xlsx", ".xls"}:
        try:
            import pandas as pd
        except ImportError as e:
            raise ImportError("pandas not installed for Excel preview") from e

        try:
            # 首先尝试使用 openpyxl 检查文件结构（需要重新创建 BytesIO）
            sheet_names = None
            preferred_sheet = None
            try:
                from openpyxl import load_workbook
                wb = load_workbook(BytesIO(content_bytes), read_only=True, data_only=True)
                sheet_names = wb.sheetnames
                logger.debug(f"Excel file has {len(sheet_names)} sheet(s): {sheet_names}")
                
                # 检查每个工作表是否有数据
                if sheet_names:
                    for sheet_name in sheet_names:
                        ws = wb[sheet_name]
                        has_data = False
                        # 检查前20行，看是否有非空数据
                        for row in ws.iter_rows(min_row=1, max_row=20, values_only=True):
                            if any(cell is not None and str(cell).strip() for cell in row):
                                has_data = True
                                break
                        
                        if has_data:
                            preferred_sheet = sheet_name
                            logger.info(f"Found data in sheet '{sheet_name}'")
                            break
                    
                    if not preferred_sheet and len(sheet_names) > 0:
                        logger.info(f"No data found in any sheet. Available sheets: {sheet_names}")
            except ImportError:
                # openpyxl 不可用，继续使用 pandas
                logger.debug("openpyxl not available, using pandas only")
            except Exception as e:
                # openpyxl 检查失败，继续使用 pandas
                logger.warning(f"openpyxl check failed: {e}, falling back to pandas")
            
            # 如果 nrows=None（即 max_rows=-1），且有多个 sheet，则读取所有 sheet
            if nrows is None and sheet_names and len(sheet_names) > 1:
                logger.info(f"Reading all {len(sheet_names)} sheets with unlimited rows")
                all_sheets_content = []
                
                for sheet_name in sheet_names:
                    try:
                        df = pd.read_excel(BytesIO(content_bytes), sheet_name=sheet_name, nrows=None, engine=None)
                        # 将 NaN 值替换为空字符串，避免在 markdown 中显示为 "nan"
                        df = df.fillna('')
                        
                        markdown_content = df.to_markdown(index=False)
                        if markdown_content and markdown_content.strip():
                            all_sheets_content.append(f"## 工作表: {sheet_name}\n\n{markdown_content}")
                    except Exception as e:
                        logger.warning(f"Failed to read sheet '{sheet_name}': {e}")
                        all_sheets_content.append(f"## 工作表: {sheet_name}\n\n*读取失败: {str(e)}*")
                
                if all_sheets_content:
                    return "\n\n---\n\n".join(all_sheets_content)
                else:
                    return "*所有工作表都为空或无数据*"
            
            # 使用 pandas 读取（优先使用有数据的表，否则使用第一个）
            target_sheet = preferred_sheet if preferred_sheet else 0
            df = pd.read_excel(BytesIO(content_bytes), sheet_name=target_sheet, nrows=nrows, engine=None)
            
            # 检查 DataFrame 是否真的为空
            if df.empty:
                # 尝试检查是否有列名但没有数据行
                if len(df.columns) > 0:
                    # 有列但无数据，可能是表头存在但数据为空
                    logger.info("Excel file has columns but no data rows")
                    return f"*文件包含列名但无数据行*\n\n列名: {', '.join(str(col) for col in df.columns)}"
                
                # 尝试读取所有工作表
                try:
                    excel_file = pd.ExcelFile(BytesIO(content_bytes))
                    all_sheets = excel_file.sheet_names
                    logger.info(f"Excel file has {len(all_sheets)} sheet(s): {all_sheets}")
                    
                    if len(all_sheets) > 1:
                        # 尝试其他工作表
                        for sheet_name in all_sheets:
                            if sheet_name == target_sheet:
                                continue
                            try:
                                df_alt = pd.read_excel(BytesIO(content_bytes), sheet_name=sheet_name, nrows=nrows)
                                if not df_alt.empty:
                                    # max_rows=-1 时不过滤空列空行，其他情况清理 DataFrame
                                    if nrows is not None:
                                        df_alt = df_alt.dropna(how='all').dropna(axis=1, how='all')
                                    
                                    if not df_alt.empty:
                                        # 将 NaN 值替换为空字符串，避免在 markdown 中显示为 "nan"
                                        df_alt = df_alt.fillna('')
                                        markdown_content = df_alt.to_markdown(index=False)
                                        if markdown_content and markdown_content.strip():
                                            return f"*工作表: {sheet_name}*\n\n{markdown_content}"
                            except Exception:
                                continue
                    
                    return f"*文件为空或没有数据*\n\n可用工作表: {', '.join(all_sheets)}"
                except Exception as e2:
                    logger.warning(f"Failed to read all sheets: {e2}")
                    return "*文件为空或没有数据*"
            
            # max_rows=-1 时不过滤空列空行，其他情况清理 DataFrame：移除完全为空的行和列
            if nrows is not None:
                df = df.dropna(how='all').dropna(axis=1, how='all')
            
            if df.empty:
                return "*文件为空或没有数据*"
            
            # 将 NaN 值替换为空字符串，避免在 markdown 中显示为 "nan"
            df = df.fillna('')
            markdown_content = df.to_markdown(index=False)
            if not markdown_content or not markdown_content.strip():
                return "*无法生成预览内容*"
            
            # 如果使用了非第一个工作表，添加提示
            if preferred_sheet and preferred_sheet != (sheet_names[0] if sheet_names else None):
                return f"*工作表: {preferred_sheet}*\n\n{markdown_content}"
            
            return markdown_content
            
        except Exception as e:
            # 记录详细错误信息用于调试
            logger.error(f"Failed to read Excel file: {type(e).__name__}: {e}", exc_info=True)
            # 如果读取失败，返回错误信息而不是空内容
            error_msg = str(e)
            if "password" in error_msg.lower() or "encrypted" in error_msg.lower():
                return "*文件受密码保护，无法预览*"
            elif "corrupt" in error_msg.lower() or "invalid" in error_msg.lower():
                return f"*文件可能已损坏: {error_msg}*"
            else:
                return f"*无法读取Excel文件: {error_msg}*"

    if ext == ".csv":
        try:
            import pandas as pd
        except ImportError as e:
            raise ImportError("pandas not installed for CSV preview") from e

        try:
            df = pd.read_csv(BytesIO(content_bytes), nrows=nrows)
        except UnicodeDecodeError:
            df = pd.read_csv(BytesIO(content_bytes), nrows=nrows, encoding="latin1")
        if df.empty:
            return "*文件为空或没有数据*"
        markdown_content = df.to_markdown(index=False)
        return markdown_content if markdown_content else "*无法生成预览内容*"

    if ext == ".docx":
        try:
            return _docx_bytes_to_markdown(content_bytes)
        except ImportError as e:
            raise ImportError("python-docx not installed for DOCX preview") from e

    if ext == ".pptx":
        try:
            from pptx import Presentation
        except ImportError as e:
            raise ImportError("python-pptx not installed for PPTX preview") from e

        prs = Presentation(BytesIO(content_bytes))
        slides = []
        for slide in prs.slides:
            texts = []
            for shape in slide.shapes:
                if hasattr(shape, "text") and shape.text:
                    text = shape.text.strip()
                    if text:
                        texts.append(text)
            if texts:
                slides.append("\n".join(texts))
        return "\n\n---\n\n".join(slides)

    if ext == ".pdf":
        try:
            from PyPDF2 import PdfReader
        except ImportError as e:
            raise ImportError("PyPDF2 not installed for PDF preview") from e

        reader = PdfReader(BytesIO(content_bytes))
        parts = []
        for page in reader.pages:
            text = page.extract_text() or ""
            text = text.strip()
            if text:
                parts.append(text)
        return "\n\n---\n\n".join(parts)

    # 未匹配的类型，尝试按文本解码
    return content_bytes.decode("utf-8", errors="replace")


async def api_workspace_file_preview(request: Request):
    """预览工作区文件。非文本文件转换为 Markdown 文本返回。
    
    Query Parameters:
        user_id: 用户ID
        chat_id: 聊天ID
        file_path: 文件路径
        max_rows: 最大读取行数（可选），不传或传正整数表示读取指定行数（默认200），-1表示读取全部
    """
    try:
        user_id = request.query_params.get("user_id") or request.headers.get("x-user-id")
        chat_id = request.query_params.get("chat_id") or request.headers.get("x-chat-id")
        file_path = request.query_params.get("file_path")
        
        # 获取 max_rows 参数，默认 200 行（向后兼容）
        max_rows_param = request.query_params.get("max_rows")
        if max_rows_param is not None:
            try:
                max_rows = int(max_rows_param)
                # -1 表示读取全部，其他小于 1 的值使用默认值
                if max_rows < 1 and max_rows != -1:
                    max_rows = 200
                # max_rows == -1 时保持 -1，传递给 _bytes_to_markdown 处理
            except (ValueError, TypeError):
                max_rows = 200  # 解析失败时使用默认值
        else:
            max_rows = 200  # 未提供参数时使用默认值

        if not user_id or not chat_id:
            return JSONResponse({
                "success": False,
                "error": "Missing required parameters: user_id and chat_id are required"
            }, status_code=400)

        if not file_path:
            return JSONResponse({
                "success": False,
                "error": "Missing required parameter: file_path"
            }, status_code=400)

        workspace_name = get_workspace_name(user_id, chat_id)
        if not workspace_name:
            return JSONResponse({
                "success": False,
                "error": "Invalid user_id and chat_id combination"
            }, status_code=400)

        user_data_dir = get_user_data_dir().resolve()
        ws_path = (user_data_dir / workspace_name).resolve()

        # Security check
        try:
            ws_path.relative_to(user_data_dir)
        except ValueError:
            return JSONResponse({
                "success": False,
                "error": "Invalid workspace path"
            }, status_code=400)

        if not ws_path.exists() or not ws_path.is_dir():
            return JSONResponse({
                "success": False,
                "error": "Workspace not found"
            }, status_code=404)

        normalized_path = file_path.lstrip("/").replace("\\", "/")
        target_file = (ws_path / normalized_path).resolve()

        # Security check: ensure file is within workspace
        try:
            target_file.relative_to(ws_path)
        except ValueError:
            return JSONResponse({
                "success": False,
                "error": "Access denied: file path outside workspace"
            }, status_code=403)

        if not target_file.exists():
            return JSONResponse({
                "success": False,
                "error": "File not found"
            }, status_code=404)

        if not target_file.is_file():
            return JSONResponse({
                "success": False,
                "error": "Path is not a file"
            }, status_code=400)

        config = load_config()
        max_file_size_mb = config.get("command", {}).get("limits", {}).get("max_file_size_mb", 100)
        max_file_size_bytes = max_file_size_mb * 1024 * 1024

        file_size = target_file.stat().st_size
        if file_size > max_file_size_bytes:
            return JSONResponse({
                "success": False,
                "error": f"File too large. Maximum size: {max_file_size_mb}MB"
            }, status_code=413)

        content_bytes = target_file.read_bytes()
        suffix = target_file.suffix.lower()

        try:
            content_md = _bytes_to_markdown(content_bytes, suffix, max_rows=max_rows)
            # 确保content_md是字符串类型且不为None
            if content_md is None:
                content_md = "*无法生成预览内容*"
            elif not isinstance(content_md, str):
                content_md = str(content_md)
            # 如果内容为空，返回提示信息
            if not content_md.strip():
                content_md = "*文件内容为空*"
        except ImportError as e:
            logger.error(f"Dependency missing for preview: {e}", exc_info=True)
            return JSONResponse({
                "success": False,
                "error": str(e)
            }, status_code=500)
        except Exception as e:
            logger.error(f"Failed to render preview for {target_file}: {e}", exc_info=True)
            return JSONResponse({
                "success": False,
                "error": f"Failed to render file: {str(e)}"
            }, status_code=500)

        return Response(content=content_md, media_type="text/markdown; charset=utf-8")

    except Exception as e:
        logger.error(f"Error in workspace file preview API: {e}", exc_info=True)
        return JSONResponse({
            "success": False,
            "error": f"Internal server error: {str(e)}"
        }, status_code=500)


async def api_workspace_image(request: Request):
    """Serve generated images from workspace/images directory."""
    try:
        user_id = request.query_params.get("user_id") or request.headers.get("x-user-id")
        chat_id = request.query_params.get("chat_id") or request.headers.get("x-chat-id")
        path = request.path_params.get("path", "")

        if not user_id or not chat_id:
            return JSONResponse({
                "success": False,
                "error": "Missing required parameters: user_id and chat_id are required"
            }, status_code=400)

        if not path:
            return JSONResponse({
                "success": False,
                "error": "Missing image path"
            }, status_code=400)

        workspace_name = get_workspace_name(user_id, chat_id)
        if not workspace_name:
            return JSONResponse({
                "success": False,
                "error": "Invalid user_id and chat_id combination"
            }, status_code=400)

        user_data_dir = get_user_data_dir().resolve()
        ws_path = (user_data_dir / workspace_name).resolve()

        # Security check
        try:
            ws_path.relative_to(user_data_dir)
        except ValueError:
            return JSONResponse({
                "success": False,
                "error": "Invalid workspace path"
            }, status_code=400)

        if not ws_path.exists() or not ws_path.is_dir():
            return JSONResponse({
                "success": False,
                "error": "Workspace not found"
            }, status_code=404)

        # Build target file path
        normalized_path = path.lstrip("/").replace("\\", "/")
        target_file = (ws_path / normalized_path).resolve()

        # Security check: ensure file is within workspace
        try:
            target_file.relative_to(ws_path)
        except ValueError:
            return JSONResponse({
                "success": False,
                "error": "Access denied: file path outside workspace"
            }, status_code=403)

        if not target_file.exists():
            return JSONResponse({
                "success": False,
                "error": "Image not found"
            }, status_code=404)

        if not target_file.is_file():
            return JSONResponse({
                "success": False,
                "error": "Path is not a file"
            }, status_code=400)

        # Determine media type based on file extension
        if target_file.suffix.lower() == '.svg':
            media_type = "image/svg+xml"
        elif target_file.suffix.lower() == '.png':
            media_type = "image/png"
        elif target_file.suffix.lower() in ['.jpg', '.jpeg']:
            media_type = "image/jpeg"
        elif target_file.suffix.lower() == '.gif':
            media_type = "image/gif"
        else:
            media_type = "application/octet-stream"

        return FileResponse(
            target_file,
            media_type=media_type,
            headers={
                "Cache-Control": "public, max-age=3600"
            }
        )

    except Exception as e:
        logger.error(f"Error serving image: {e}", exc_info=True)
        return JSONResponse({
            "success": False,
            "error": f"Internal server error: {str(e)}"
        }, status_code=500)


async def api_public_image(request: Request):
    """公开访问生成的图片（无需 user_id 和 chat_id）。
    
    从公共图床目录读取图片，路径格式：MCP_WORKSPACES_DIR/YYYYMMDD/filename
    支持按日期目录组织，便于管理和清理。
    """
    try:
        filename = request.path_params.get("filename", "")
        
        if not filename:
            return JSONResponse({
                "success": False,
                "error": "Missing image filename"
            }, status_code=400)
        
        # 安全检查：只允许图片文件
        allowed_extensions = {'.png', '.svg', '.jpg', '.jpeg', '.gif', '.webp'}
        file_ext = Path(filename).suffix.lower()
        if file_ext not in allowed_extensions:
            return JSONResponse({
                "success": False,
                "error": "Invalid file type"
            }, status_code=400)
        
        # 安全检查：文件名必须包含 UUID 格式（至少 8 个十六进制字符）
        import re
        if not re.search(r'[0-9a-f]{8,}', filename, re.IGNORECASE):
            return JSONResponse({
                "success": False,
                "error": "Invalid filename format"
            }, status_code=400)
        
        user_data_dir = get_user_data_dir().resolve()
        
        if not user_data_dir.exists():
            return JSONResponse({
                "success": False,
                "error": "Image not found"
            }, status_code=404)
        
        # 从图床目录查找文件（按日期目录组织）
        # 优先查找今天的目录，如果不存在则查找最近30天的目录
        from datetime import datetime, timedelta
        target_file = None
        
        # 首先尝试今天的日期目录
        today_str = datetime.now().strftime("%Y%m%d")
        today_dir = user_data_dir / today_str
        if today_dir.exists() and today_dir.is_dir():
            candidate_file = today_dir / filename
            if candidate_file.exists() and candidate_file.is_file():
                try:
                    candidate_file.resolve().relative_to(user_data_dir)
                    target_file = candidate_file
                except ValueError:
                    pass
        
        # 如果今天目录没找到，查找最近30天的目录
        if target_file is None:
            for days_ago in range(1, 31):
                date_str = (datetime.now() - timedelta(days=days_ago)).strftime("%Y%m%d")
                date_dir = user_data_dir / date_str
                if date_dir.exists() and date_dir.is_dir():
                    candidate_file = date_dir / filename
                    if candidate_file.exists() and candidate_file.is_file():
                        try:
                            candidate_file.resolve().relative_to(user_data_dir)
                            target_file = candidate_file
                            break
                        except ValueError:
                            continue
        
        if target_file is None or not target_file.exists():
            return JSONResponse({
                "success": False,
                "error": "Image not found"
            }, status_code=404)
        
        # 确定媒体类型
        if file_ext == '.svg':
            media_type = "image/svg+xml"
        elif file_ext == '.png':
            media_type = "image/png"
        elif file_ext in ['.jpg', '.jpeg']:
            media_type = "image/jpeg"
        elif file_ext == '.gif':
            media_type = "image/gif"
        elif file_ext == '.webp':
            media_type = "image/webp"
        else:
            media_type = "application/octet-stream"
        
        return FileResponse(
            target_file,
            media_type=media_type,
            headers={
                "Cache-Control": "public, max-age=3600"
            }
        )
        
    except Exception as e:
        logger.error(f"Error serving public image: {e}", exc_info=True)
        return JSONResponse({
            "success": False,
            "error": f"Internal server error: {str(e)}"
        }, status_code=500)


async def api_workspace_download(request: Request):
    """Download entire workspace as a zip file."""
    temp_zip_path = None
    try:
        user_id = request.query_params.get("user_id") or request.headers.get("x-user-id")
        chat_id = request.query_params.get("chat_id") or request.headers.get("x-chat-id")

        if not user_id or not chat_id:
            return JSONResponse({
                "success": False,
                "error": "Missing required parameters: user_id and chat_id are required"
            }, status_code=400)

        workspace_name = get_workspace_name(user_id, chat_id)
        if not workspace_name:
            return JSONResponse({
                "success": False,
                "error": "Invalid user_id and chat_id combination"
            }, status_code=400)

        user_data_dir = get_user_data_dir().resolve()
        ws_path = (user_data_dir / workspace_name).resolve()

        # Security check
        try:
            ws_path.relative_to(user_data_dir)
        except ValueError:
            return JSONResponse({
                "success": False,
                "error": "Invalid workspace path"
            }, status_code=400)

        if not ws_path.exists() or not ws_path.is_dir():
            return JSONResponse({
                "success": False,
                "error": "Workspace not found"
            }, status_code=404)

        # Check workspace size limit
        config = load_config()
        max_file_size_mb = config.get("command", {}).get("limits", {}).get("max_file_size_mb", 100)
        max_file_size_bytes = max_file_size_mb * 1024 * 1024

        # Calculate total size of workspace
        total_size = 0
        for root, dirs, files in os.walk(ws_path):
            for file in files:
                file_path = Path(root) / file
                try:
                    total_size += file_path.stat().st_size
                    if total_size > max_file_size_bytes:
                        return JSONResponse({
                            "success": False,
                            "error": f"Workspace too large. Maximum size: {max_file_size_mb}MB"
                        }, status_code=413)
                except (OSError, PermissionError):
                    continue

        # Create temporary zip file
        temp_dir = tempfile.gettempdir()
        zip_filename = f"workspace_{workspace_name}.zip"
        temp_zip_path = Path(temp_dir) / zip_filename

        # Create zip file
        try:
            with zipfile.ZipFile(temp_zip_path, 'w', zipfile.ZIP_DEFLATED) as zipf:
                for root, dirs, files in os.walk(ws_path):
                    for file in files:
                        file_path = Path(root) / file
                        try:
                            arcname = file_path.relative_to(ws_path)
                            zipf.write(file_path, arcname)
                        except (OSError, PermissionError) as e:
                            logger.warning(f"Skipping file {file_path}: {e}")
                            continue

            # Check if zip file was created successfully
            if not temp_zip_path.exists() or temp_zip_path.stat().st_size == 0:
                if temp_zip_path.exists():
                    temp_zip_path.unlink()
                return JSONResponse({
                    "success": False,
                    "error": "Failed to create zip file or workspace is empty"
                }, status_code=500)
        except Exception as zip_error:
            if temp_zip_path and temp_zip_path.exists():
                try:
                    temp_zip_path.unlink()
                except Exception:
                    pass
            logger.error(f"Error creating zip file: {zip_error}", exc_info=True)
            return JSONResponse({
                "success": False,
                "error": f"Failed to create zip file: {str(zip_error)}"
            }, status_code=500)

        # Stream zip file and clean up after download
        def generate():
            try:
                with open(temp_zip_path, 'rb') as f:
                    while True:
                        chunk = f.read(8192)  # 8KB chunks
                        if not chunk:
                            break
                        yield chunk
            finally:
                # Clean up temp file after streaming
                try:
                    if temp_zip_path.exists():
                        temp_zip_path.unlink()
                        logger.debug(f"Cleaned up temporary file: {temp_zip_path}")
                except Exception as e:
                    logger.warning(f"Failed to clean up temporary file {temp_zip_path}: {e}")

        return StreamingResponse(
            generate(),
            media_type="application/zip",
            headers={
                "Content-Disposition": f'attachment; filename="{zip_filename}"'
            }
        )

    except Exception as e:
        logger.error(f"Error in workspace download API: {e}", exc_info=True)
        if temp_zip_path and temp_zip_path.exists():
            try:
                temp_zip_path.unlink()
            except Exception:
                pass
        return JSONResponse({
            "success": False,
            "error": f"Internal server error: {str(e)}"
        }, status_code=500)


async def api_workspace_file_upload(request: Request):
    """Upload files to workspace using multipart/form-data.

    支持三种方式：
    - 传统文件字段：`files`（可多选）或 `file`
    - 远程 URL 下载：`file_url` 或 `url`
    - 文本+文件名：`text`（文本内容）和 `filename`（文件名）
    """
    try:
        # Get user_id and chat_id from query params or headers
        user_id = request.query_params.get("user_id") or request.headers.get("x-user-id")
        chat_id = request.query_params.get("chat_id") or request.headers.get("x-chat-id")
        
        # Get target directory from form data or query params (optional, defaults to root)
        target_dir = request.query_params.get("target_dir", "")

        if not user_id or not chat_id:
            return JSONResponse({
                "success": False,
                "error": "Missing required parameters: user_id and chat_id are required"
            }, status_code=400)

        workspace_name = get_workspace_name(user_id, chat_id)
        if not workspace_name:
            return JSONResponse({
                "success": False,
                "error": "Invalid user_id and chat_id combination"
            }, status_code=400)

        user_data_dir = get_user_data_dir().resolve()
        ws_path = (user_data_dir / workspace_name).resolve()

        # Security check
        try:
            ws_path.relative_to(user_data_dir)
        except ValueError:
            return JSONResponse({
                "success": False,
                "error": "Invalid workspace path"
            }, status_code=400)

        # Create workspace if it doesn't exist
        ws_path.mkdir(parents=True, exist_ok=True)

        # Parse multipart form data
        form = await request.form()
        
        # Check for text + filename upload (new method)
        text_content = form.get("text")
        filename_from_text = form.get("filename")
        
        # Get files from form (support multiple files with same field name)
        files = form.getlist("files")
        if not files:
            # Also try "file" as singular field name
            files = form.getlist("file")

        # Also support remote URLs
        file_urls = form.getlist("file_url")
        if not file_urls:
            file_urls = form.getlist("url")
        
        if not files and not file_urls and not (text_content and filename_from_text):
            return JSONResponse({
                "success": False,
                "error": "No files provided. Use 'files'/'file', provide 'file_url'/'url', or use 'text' + 'filename'."
            }, status_code=400)

        # Get file size limit from config
        config = load_config()
        max_file_size_mb = config.get("command", {}).get("limits", {}).get("max_file_size_mb", 100)
        max_file_size_bytes = max_file_size_mb * 1024 * 1024

        # Determine target directory
        target_path = ws_path
        if target_dir:
            # Normalize target directory path
            normalized_dir = target_dir.lstrip("/").replace("\\", "/")
            target_path = (ws_path / normalized_dir).resolve()
            
            # Security check: ensure target directory is within workspace
            try:
                target_path.relative_to(ws_path)
            except ValueError:
                return JSONResponse({
                    "success": False,
                    "error": "Access denied: target directory outside workspace"
                }, status_code=403)
            
            # Create target directory if it doesn't exist
            target_path.mkdir(parents=True, exist_ok=True)

        uploaded_files = []
        errors = []

        # Process text + filename upload (if provided)
        if text_content is not None and filename_from_text:
            try:
                # Validate filename
                if not isinstance(filename_from_text, str) or not filename_from_text.strip():
                    errors.append({
                        "file": "text+filename",
                        "error": "Invalid filename"
                    })
                else:
                    # Security: sanitize filename to prevent path traversal
                    safe_filename = os.path.basename(filename_from_text).replace("..", "").replace("/", "").replace("\\", "")
                    if not safe_filename:
                        errors.append({
                            "file": filename_from_text,
                            "error": "Invalid filename"
                        })
                    else:
                        # Convert text to bytes (UTF-8 encoding)
                        if isinstance(text_content, str):
                            content_bytes = text_content.encode("utf-8")
                        else:
                            # If it's already bytes, use as-is
                            content_bytes = text_content if isinstance(text_content, bytes) else str(text_content).encode("utf-8")
                        
                        file_size = len(content_bytes)
                        
                        # Check file size limit
                        if file_size > max_file_size_bytes:
                            errors.append({
                                "file": safe_filename,
                                "error": f"File too large. Maximum size: {max_file_size_mb}MB"
                            })
                        else:
                            # Build target file path
                            target_file = target_path / safe_filename
                            
                            # Check if file already exists
                            overwrite = target_file.exists()
                            
                            # Write file
                            target_file.parent.mkdir(parents=True, exist_ok=True)
                            with open(target_file, "wb") as f:
                                f.write(content_bytes)
                            
                            # Get relative path from workspace root
                            try:
                                relative_path = str(target_file.relative_to(ws_path))
                            except ValueError:
                                relative_path = safe_filename
                            
                            uploaded_files.append({
                                "filename": safe_filename,
                                "original_filename": filename_from_text,
                                "path": relative_path,
                                "size": file_size,
                                "size_human": human_size(file_size),
                                "overwritten": overwrite,
                                "source": "text"
                            })
                            
                            logger.info(f"Uploaded file from text: {target_file} ({file_size} bytes)")
                            
            except Exception as e:
                logger.error(f"Error uploading file from text: {e}", exc_info=True)
                errors.append({
                    "file": filename_from_text if filename_from_text else "text+filename",
                    "error": str(e)
                })

        # Process each file
        for file_item in files:
            try:
                # Check if it's a file upload
                if not hasattr(file_item, "filename") or not file_item.filename:
                    errors.append({
                        "file": "unknown",
                        "error": "Invalid file item"
                    })
                    continue

                # Some clients (notably on Windows) may URL-encode the filename (e.g. spaces -> %20).
                # Decode it to ensure we persist the real filename on disk.
                filename = unquote(file_item.filename)
                
                # Security: sanitize filename to prevent path traversal
                # Remove any path separators and dangerous characters
                safe_filename = os.path.basename(filename).replace("..", "").replace("/", "").replace("\\", "")
                if not safe_filename:
                    errors.append({
                        "file": filename,
                        "error": "Invalid filename"
                    })
                    continue

                # Build target file path
                target_file = target_path / safe_filename

                # Read file content
                content = await file_item.read()
                file_size = len(content)

                # Check file size limit
                if file_size > max_file_size_bytes:
                    errors.append({
                        "file": filename,
                        "error": f"File too large. Maximum size: {max_file_size_mb}MB"
                    })
                    continue

                # Check if file already exists
                overwrite = target_file.exists()

                # Write file
                target_file.parent.mkdir(parents=True, exist_ok=True)
                with open(target_file, "wb") as f:
                    f.write(content)

                # Get relative path from workspace root
                try:
                    relative_path = str(target_file.relative_to(ws_path))
                except ValueError:
                    relative_path = safe_filename

                uploaded_files.append({
                    "filename": safe_filename,
                    "original_filename": filename,
                    "path": relative_path,
                    "size": file_size,
                    "size_human": human_size(file_size),
                    "overwritten": overwrite
                })

                logger.info(f"Uploaded file: {target_file} ({file_size} bytes)")

            except Exception as e:
                logger.error(f"Error uploading file {file_item.filename if hasattr(file_item, 'filename') else 'unknown'}: {e}", exc_info=True)
                errors.append({
                    "file": file_item.filename if hasattr(file_item, "filename") else "unknown",
                    "error": str(e)
                })

        # Process each remote URL (download then save)
        if file_urls:
            async with httpx.AsyncClient(timeout=60.0, follow_redirects=True) as client:
                for url in file_urls:
                    try:
                        if not isinstance(url, str) or not url.strip():
                            errors.append({
                                "file": "url",
                                "error": "Empty url value"
                            })
                            continue

                        url = url.strip()
                        parsed = urlparse(url)
                        if parsed.scheme not in ("http", "https"):
                            errors.append({
                                "file": url,
                                "error": "Only http/https URLs are supported"
                            })
                            continue

                        try:
                            resp = await client.get(url)
                        except httpx.HTTPError as e:
                            errors.append({
                                "file": url,
                                "error": f"Failed to download: {e}"
                            })
                            continue

                        status_code = resp.status_code
                        if status_code >= 400:
                            errors.append({
                                "file": url,
                                "error": f"HTTP {status_code}"
                            })
                            continue

                        content = resp.content or b""
                        file_size = len(content)

                        # Check file size limit
                        if file_size > max_file_size_bytes:
                            errors.append({
                                "file": url,
                                "error": f"File too large. Maximum size: {max_file_size_mb}MB"
                            })
                            continue

                        # Determine filename from URL or content-type
                        filename = unquote(os.path.basename(parsed.path))
                        if not filename:
                            ctype = resp.headers.get("content-type", "")
                            ext = mimetypes.guess_extension(ctype.split(";")[0].strip()) if ctype else None
                            filename = f"downloaded_file{ext or ''}"

                        safe_filename = os.path.basename(filename).replace("..", "").replace("/", "").replace("\\", "")
                        if not safe_filename:
                            safe_filename = "downloaded_file"

                        target_file = target_path / safe_filename
                        overwrite = target_file.exists()

                        target_file.parent.mkdir(parents=True, exist_ok=True)
                        with open(target_file, "wb") as f:
                            f.write(content)

                        try:
                            relative_path = str(target_file.relative_to(ws_path))
                        except ValueError:
                            relative_path = safe_filename

                        uploaded_files.append({
                            "filename": safe_filename,
                            "original_filename": safe_filename,
                            "path": relative_path,
                            "size": file_size,
                            "size_human": human_size(file_size),
                            "overwritten": overwrite,
                            "source_url": url
                        })

                        logger.info(f"Downloaded file from URL {url} -> {target_file} ({file_size} bytes)")

                    except Exception as e:
                        logger.error(f"Error downloading file from url {url}: {e}", exc_info=True)
                        errors.append({
                            "file": url,
                            "error": str(e)
                        })

        # Return response
        if not uploaded_files and errors:
            return JSONResponse({
                "success": False,
                "error": "All files failed to upload",
                "errors": errors
            }, status_code=400)

        return JSONResponse({
            "success": True,
            "uploaded_count": len(uploaded_files),
            "files": uploaded_files,
            "errors": errors if errors else None,
            "target_dir": str(target_path.relative_to(ws_path)) if target_path != ws_path else "/"
        })

    except Exception as e:
        logger.error(f"Error in workspace file upload API: {e}", exc_info=True)
        return JSONResponse({
            "success": False,
            "error": f"Internal server error: {str(e)}"
        }, status_code=500)


async def api_workspace_file_update(request: Request):
    """Update or create a text-editable file in the workspace via POST."""
    try:
        try:
            body = await request.json()
        except Exception:
            return JSONResponse({
                "success": False,
                "error": "Invalid JSON body"
            }, status_code=400)

        user_id = body.get("user_id") or request.query_params.get("user_id") or request.headers.get("x-user-id")
        chat_id = body.get("chat_id") or request.query_params.get("chat_id") or request.headers.get("x-chat-id")
        file_path = body.get("file_path") or body.get("path")
        content = body.get("content")

        if not user_id or not chat_id:
            return JSONResponse({
                "success": False,
                "error": "Missing required parameters: user_id and chat_id are required"
            }, status_code=400)

        if not file_path:
            return JSONResponse({
                "success": False,
                "error": "Missing required parameter: file_path"
            }, status_code=400)

        if not isinstance(content, str):
            return JSONResponse({
                "success": False,
                "error": "Content must be a string"
            }, status_code=400)

        workspace_name = get_workspace_name(user_id, chat_id)
        if not workspace_name:
            return JSONResponse({
                "success": False,
                "error": "Invalid user_id and chat_id combination"
            }, status_code=400)

        user_data_dir = get_user_data_dir().resolve()
        ws_path = (user_data_dir / workspace_name).resolve()

        try:
            ws_path.relative_to(user_data_dir)
        except ValueError:
            return JSONResponse({
                "success": False,
                "error": "Invalid workspace path"
            }, status_code=400)

        ws_path.mkdir(parents=True, exist_ok=True)

        normalized_path = file_path.lstrip("/").replace("\\", "/")
        target_file = (ws_path / normalized_path).resolve()

        try:
            target_file.relative_to(ws_path)
        except ValueError:
            return JSONResponse({
                "success": False,
                "error": "Access denied: file path outside workspace"
            }, status_code=403)

        if target_file.exists() and target_file.is_dir():
            return JSONResponse({
                "success": False,
                "error": "Path is a directory, not a file"
            }, status_code=400)

        allowed_ext = {
            ".txt", ".md", ".markdown", ".json", ".yaml", ".yml", ".html", ".htm",
            ".css", ".js", ".ts", ".tsx", ".py", ".ini", ".toml", ".env", ".xml",
            ".csv"
        }
        ext = target_file.suffix.lower()
        if ext and ext not in allowed_ext:
            return JSONResponse({
                "success": False,
                "error": f"File type not supported for text edit: {ext}"
            }, status_code=400)

        # Enforce size limit using the same config as uploads
        config = load_config()
        max_file_size_mb = config.get("command", {}).get("limits", {}).get("max_file_size_mb", 100)
        max_file_size_bytes = max_file_size_mb * 1024 * 1024
        content_bytes = content.encode("utf-8")
        if len(content_bytes) > max_file_size_bytes:
            return JSONResponse({
                "success": False,
                "error": f"Content too large. Maximum size: {max_file_size_mb}MB"
            }, status_code=413)

        target_file.parent.mkdir(parents=True, exist_ok=True)
        overwritten = target_file.exists()
        with open(target_file, "w", encoding="utf-8") as f:
            f.write(content)

        try:
            relative_path = str(target_file.relative_to(ws_path))
        except ValueError:
            relative_path = target_file.name

        size_written = target_file.stat().st_size

        return JSONResponse({
            "success": True,
            "path": relative_path,
            "size": size_written,
            "size_human": human_size(size_written),
            "overwritten": overwritten,
        })

    except Exception as e:
        logger.error(f"Error in workspace file update API: {e}", exc_info=True)
        return JSONResponse({
            "success": False,
            "error": f"Internal server error: {str(e)}"
        }, status_code=500)


async def api_workspace_file_delete(request: Request):
    """Delete a specific file in the workspace via POST."""
    try:
        try:
            body = await request.json()
        except Exception:
            return JSONResponse({
                "success": False,
                "error": "Invalid JSON body"
            }, status_code=400)

        user_id = body.get("user_id") or request.query_params.get("user_id") or request.headers.get("x-user-id")
        chat_id = body.get("chat_id") or request.query_params.get("chat_id") or request.headers.get("x-chat-id")
        file_path = body.get("file_path") or body.get("path")

        if not user_id or not chat_id:
            return JSONResponse({
                "success": False,
                "error": "Missing required parameters: user_id and chat_id are required"
            }, status_code=400)

        if not file_path:
            return JSONResponse({
                "success": False,
                "error": "Missing required parameter: file_path"
            }, status_code=400)

        workspace_name = get_workspace_name(user_id, chat_id)
        if not workspace_name:
            return JSONResponse({
                "success": False,
                "error": "Invalid user_id and chat_id combination"
            }, status_code=400)

        user_data_dir = get_user_data_dir().resolve()
        ws_path = (user_data_dir / workspace_name).resolve()

        try:
            ws_path.relative_to(user_data_dir)
        except ValueError:
            return JSONResponse({
                "success": False,
                "error": "Invalid workspace path"
            }, status_code=400)

        normalized_path = file_path.lstrip("/").replace("\\", "/")
        target_file = (ws_path / normalized_path).resolve()

        try:
            target_file.relative_to(ws_path)
        except ValueError:
            return JSONResponse({
                "success": False,
                "error": "Access denied: file path outside workspace"
            }, status_code=403)

        if not target_file.exists():
            return JSONResponse({
                "success": False,
                "error": "File not found"
            }, status_code=404)

        if not target_file.is_file():
            return JSONResponse({
                "success": False,
                "error": "Path is not a file"
            }, status_code=400)

        size_bytes = target_file.stat().st_size
        target_file.unlink()

        return JSONResponse({
            "success": True,
            "message": "File deleted",
            "path": str(target_file.relative_to(ws_path)),
            "size": size_bytes,
            "size_human": human_size(size_bytes),
        })

    except Exception as e:
        logger.error(f"Error in workspace file delete API: {e}", exc_info=True)
        return JSONResponse({
            "success": False,
            "error": f"Internal server error: {str(e)}"
        }, status_code=500)


# ========== Preview Deployment API ==========


def get_components_cache():
    """Get the global components cache (avoids circular import)."""
    # Import here to avoid circular import with server.py
    from . import server
    return server._components_cache


async def api_preview_deploy(request: Request):
    """Deploy/restart preview for a workspace."""
    try:
        # Get parameters from query string or JSON body
        if request.method == "POST":
            try:
                body = await request.json()
                user_id = body.get("user_id") or request.query_params.get("user_id")
                chat_id = body.get("chat_id") or request.query_params.get("chat_id")
                directory = body.get("directory", "/")
                entry = body.get("entry", "index.html")
            except:
                user_id = request.query_params.get("user_id")
                chat_id = request.query_params.get("chat_id")
                directory = request.query_params.get("directory", "/")
                entry = request.query_params.get("entry", "index.html")
        else:
            user_id = request.query_params.get("user_id")
            chat_id = request.query_params.get("chat_id")
            directory = request.query_params.get("directory", "/")
            entry = request.query_params.get("entry", "index.html")

        if not user_id or not chat_id:
            return JSONResponse({
                "success": False,
                "error": "Missing required parameters: user_id and chat_id are required"
            }, status_code=400)

        workspace_name = get_workspace_name(user_id, chat_id)
        if not workspace_name:
            return JSONResponse({
                "success": False,
                "error": "Invalid user_id and chat_id combination"
            }, status_code=400)

        # Get or create workspace path
        user_data_dir = get_user_data_dir()
        workspace_path = user_data_dir / workspace_name
        workspace_path.mkdir(parents=True, exist_ok=True)

        # Get preview manager from cache
        _components_cache = get_components_cache()
        if workspace_name in _components_cache:
            components = _components_cache[workspace_name]
        else:
            # Create components for this workspace
            from .security import PathValidator

            validator = PathValidator([workspace_path], virtual_root=workspace_path)
            config = load_config()

            preview_manager = PreviewManager(
                workspace_path=workspace_path,
                workspace_name=workspace_name,
                config=config,
            )

            components = {
                "preview": preview_manager,
            }
            _components_cache[workspace_name] = components

        preview_mgr = components.get("preview")
        if not preview_mgr:
            return JSONResponse({
                "success": False,
                "error": "Preview manager not available"
            }, status_code=500)

        # Start or restart preview
        result = await preview_mgr.start_preview(directory, entry)

        if result.get("success"):
            return JSONResponse({
                "success": True,
                "url": result.get("url"),
                "port": result.get("port"),
                "subdomain": result.get("subdomain"),
                "directory": result.get("directory"),
                "entry": result.get("entry"),
                "message": result.get("message", "Preview deployed successfully")
            })
        else:
            return JSONResponse({
                "success": False,
                "error": result.get("error", "Failed to start preview")
            }, status_code=500)

    except Exception as e:
        logger.error(f"Error in preview deploy API: {e}", exc_info=True)
        return JSONResponse({
            "success": False,
            "error": f"Internal server error: {str(e)}"
        }, status_code=500)


async def api_preview_status(request: Request):
    """Get preview status for a workspace."""
    try:
        user_id = request.query_params.get("user_id")
        chat_id = request.query_params.get("chat_id")

        if not user_id or not chat_id:
            return JSONResponse({
                "success": False,
                "error": "Missing required parameters: user_id and chat_id are required"
            }, status_code=400)

        workspace_name = get_workspace_name(user_id, chat_id)
        if not workspace_name:
            return JSONResponse({
                "success": False,
                "error": "Invalid user_id and chat_id combination"
            }, status_code=400)

        # Get components from cache
        _components_cache = get_components_cache()
        if workspace_name not in _components_cache:
            return JSONResponse({
                "success": True,
                "active": False,
                "message": "No preview has been started for this workspace"
            })

        components = _components_cache[workspace_name]
        preview_mgr = components.get("preview")

        if not preview_mgr:
            return JSONResponse({
                "success": True,
                "active": False,
                "message": "Preview manager not available"
            })

        status = preview_mgr.get_status()
        return JSONResponse({
            "success": True,
            **status
        })

    except Exception as e:
        logger.error(f"Error in preview status API: {e}", exc_info=True)
        return JSONResponse({
            "success": False,
            "error": f"Internal server error: {str(e)}"
        }, status_code=500)


async def api_preview_stop(request: Request):
    """Stop preview for a workspace."""
    try:
        user_id = request.query_params.get("user_id")
        chat_id = request.query_params.get("chat_id")

        if not user_id or not chat_id:
            return JSONResponse({
                "success": False,
                "error": "Missing required parameters: user_id and chat_id are required"
            }, status_code=400)

        workspace_name = get_workspace_name(user_id, chat_id)
        if not workspace_name:
            return JSONResponse({
                "success": False,
                "error": "Invalid user_id and chat_id combination"
            }, status_code=400)

        # Get components from cache
        _components_cache = get_components_cache()
        if workspace_name not in _components_cache:
            return JSONResponse({
                "success": True,
                "message": "No active preview to stop"
            })

        components = _components_cache[workspace_name]
        preview_mgr = components.get("preview")

        if not preview_mgr:
            return JSONResponse({
                "success": True,
                "message": "Preview manager not available"
            })

        result = await preview_mgr.stop_preview()
        return JSONResponse({
            "success": True,
            **result
        })

    except Exception as e:
        logger.error(f"Error in preview stop API: {e}", exc_info=True)
        return JSONResponse({
            "success": False,
            "error": f"Internal server error: {str(e)}"
        }, status_code=500)


async def api_preview_access(request: Request):
    """Record that a preview URL was accessed (for LRU tracking)."""
    try:
        from .command.preview import get_global_usage_tracker

        subdomain = request.query_params.get("subdomain")
        port = request.query_params.get("port")
        url = request.query_params.get("url")

        # If URL is provided, try to extract subdomain or port
        if url and not subdomain and not port:
            import re

            config = load_config()
            preview_config = config.get("preview", {})
            wildcard_domain = preview_config.get("wildcard_domain")

            if wildcard_domain:
                domain_part = wildcard_domain.replace("https://", "").replace("http://", "").replace("*", "")
                if domain_part.startswith("."):
                    domain_part = domain_part[1:]

                domain_escaped = re.escape(domain_part)
                pattern = rf'https?://([^.]+)\.{domain_escaped}'
                subdomain_match = re.search(pattern, url)
                if subdomain_match:
                    subdomain = subdomain_match.group(1)

            if not subdomain:
                port_match = re.search(r':(\d+)/', url)
                if port_match:
                    port = int(port_match.group(1))

        if not subdomain and not port:
            return JSONResponse({
                "success": False,
                "error": "Missing required parameter: subdomain or port (or url)"
            }, status_code=400)

        usage_tracker = get_global_usage_tracker()
        usage_tracker.record_access(
            subdomain=subdomain,
            port=int(port) if port else None
        )

        return JSONResponse({
            "success": True,
            "message": "Access recorded"
        })

    except Exception as e:
        logger.error(f"Error in preview access API: {e}", exc_info=True)
        return JSONResponse({
            "success": False,
            "error": f"Internal server error: {str(e)}"
        }, status_code=500)


async def api_preview_mapping(request: Request):
    """Get subdomain to preview mapping for nginx/diagnostics."""
    try:
        from .command.preview import get_global_usage_tracker

        usage_tracker = get_global_usage_tracker()
        subdomain = request.query_params.get("subdomain")

        if subdomain:
            # Get specific subdomain mapping
            usage = usage_tracker.get_preview_by_key(subdomain)
            if usage:
                return JSONResponse({
                    "success": True,
                    "subdomain": subdomain,
                    "directory": usage.directory,
                    "entry": usage.entry,
                    "url": usage.url
                })
            else:
                return JSONResponse({
                    "success": False,
                    "error": "Subdomain not found"
                }, status_code=404)
        else:
            # Get all active mappings
            all_previews = usage_tracker.get_all_previews()
            mappings = {}
            for preview in all_previews:
                if preview.subdomain:
                    mappings[preview.subdomain] = {
                        "directory": preview.directory,
                        "entry": preview.entry,
                        "url": preview.url,
                        "workspace_name": preview.workspace_name
                    }

            return JSONResponse({
                "success": True,
                "mappings": mappings,
                "count": len(mappings)
            })

    except Exception as e:
        logger.error(f"Error in preview mapping API: {e}", exc_info=True)
        return JSONResponse({
            "success": False,
            "error": f"Internal server error: {str(e)}"
        }, status_code=500)


# ========== Python Execution API ==========

async def api_exec_python(request: Request):
    """执行 Python 代码并返回执行结果（用户侧接口）。
    
    支持 POST 请求，参数通过 JSON body 传递。
    
    参数：
    - user_id: 用户ID（必需，可通过 body 或 header x-user-id 传递）
    - chat_id: 聊天ID（必需，可通过 body 或 header x-chat-id 传递）
    - code: Python 代码字符串（必需）
    
    返回：
    {
        "success": bool,
        "exit_code": int,
        "stdout": str,
        "stderr": str,
        "execution_time": float,
        "timed_out": bool,
        "error": str (可选)
    }
    """
    try:
        # 获取参数（仅支持 POST，从 JSON body 或 header 获取）
        try:
            body = await request.json()
            user_id = body.get("user_id") or request.headers.get("x-user-id")
            chat_id = body.get("chat_id") or request.headers.get("x-chat-id")
            code = body.get("code", "")
        except:
            # 如果 JSON 解析失败，尝试从 header 获取
            user_id = request.headers.get("x-user-id")
            chat_id = request.headers.get("x-chat-id")
            code = ""

        # 验证必需参数
        if not user_id or not chat_id:
            return JSONResponse({
                "success": False,
                "error": "Missing required parameters: user_id and chat_id are required"
            }, status_code=400)

        if not code:
            return JSONResponse({
                "success": False,
                "error": "Missing required parameter: code"
            }, status_code=400)

        # 获取工作空间名称
        workspace_name = get_workspace_name(user_id, chat_id)
        if not workspace_name:
            return JSONResponse({
                "success": False,
                "error": "Invalid user_id and chat_id combination"
            }, status_code=400)

        # 获取或创建 components
        _components_cache = get_components_cache()
        
        if workspace_name in _components_cache:
            components = _components_cache[workspace_name]
        else:
            # 创建 components
            user_data_dir = get_user_data_dir()
            workspace_path = user_data_dir / workspace_name
            workspace_path.mkdir(parents=True, exist_ok=True)

            from .security import PathValidator
            from .operations import FileOperations
            from .command import CommandExecutor
            from .excel import ExcelOperations
            from .grep import GrepTools
            from .advanced import AdvancedFileOperations

            config = load_config()
            validator = PathValidator([workspace_path], virtual_root=workspace_path)
            operations = FileOperations(validator)
            advanced = AdvancedFileOperations(validator, operations)
            grep = GrepTools(validator)
            excel = ExcelOperations(validator, config.get("excel", {}))
            
            command_config = config.get("command", {})
            command_executor = CommandExecutor(
                workspace_path=workspace_path,
                config=command_config,
            )

            components = {
                "validator": validator,
                "operations": operations,
                "advanced": advanced,
                "grep": grep,
                "excel": excel,
                "command": command_executor,
                "allowed_dirs": validator.get_allowed_dirs(),
                "workspace_path": workspace_path,
            }
            _components_cache[workspace_name] = components

        # 检查 command_executor 是否可用
        command_executor = components.get("command")
        if not command_executor:
            return JSONResponse({
                "success": False,
                "error": "Command executor not available"
            }, status_code=500)

        # 执行 Python 代码
        from .tools.exec_tools import exec_command

        result = await exec_command(
            runtime="python",
            command_executor=command_executor,
            validator=components["validator"],
            operations=components["operations"],
            code=code,
            file=None,
            args=None,
        )

        return JSONResponse({
            "success": result.get("success", False),
            "exit_code": result.get("exit_code", -1),
            "stdout": result.get("stdout", ""),
            "stderr": result.get("stderr", ""),
            "execution_time": result.get("execution_time", 0),
            "timed_out": result.get("timed_out", False),
            "error": result.get("error"),
        })

    except Exception as e:
        logger.error(f"Error in exec python API: {e}", exc_info=True)
        return JSONResponse({
            "success": False,
            "error": f"Internal server error: {str(e)}"
        }, status_code=500)


# ========== Web UI Endpoints ==========

def web_disabled_response():
    """Return error response when web UI is disabled."""
    return JSONResponse({"success": False, "error": "Web UI is disabled"}, status_code=403)


def get_web_password(request: Request) -> str:
    """Get web password from request headers or query params."""
    return request.headers.get("x-web-password") or request.query_params.get("password", "")


def verify_web_auth(request: Request, web_cfg: Dict[str, Any]) -> Optional[JSONResponse]:
    """Verify web UI authentication."""
    if not web_cfg.get("enabled", False):
        return web_disabled_response()

    expected = str(web_cfg.get("password", ""))
    provided = get_web_password(request)

    if expected and provided != expected:
        return JSONResponse(
            {"success": False, "error": "Invalid password"},
            status_code=401,
        )
    return None


async def web_tree_endpoint(request: Request):
    """Web UI endpoint for file tree."""
    web_cfg = get_web_config()
    auth_error = verify_web_auth(request, web_cfg)
    if auth_error:
        return auth_error

    user_data_dir = get_user_data_dir().resolve()
    req_path = request.query_params.get("path", "/") or "/"
    try:
        max_depth = int(request.query_params.get("max_depth", "5"))
    except ValueError:
        max_depth = 5

    rel_path = req_path.lstrip("/")
    target = (user_data_dir / rel_path).resolve() if rel_path else user_data_dir

    if not str(target).startswith(str(user_data_dir)):
        return JSONResponse({"success": False, "error": "Access denied"}, status_code=403)
    if not target.exists():
        return JSONResponse({"success": False, "error": "Path not found"}, status_code=404)

    def build_tree(path: Path, rel: str = "", depth: int = 0):
        try:
            stat = path.stat()
            created = format_timestamp(stat.st_ctime)
            modified = format_timestamp(stat.st_mtime)
            size = stat.st_size
        except Exception:
            stat = None
            created = modified = None
            size = 0

        rel_clean = rel.strip("/")
        display_path = "/" if not rel_clean else f"/{rel_clean}"
        node = {
            "name": "/" if not rel_clean else path.name,
            "path": display_path,
            "type": "directory" if path.is_dir() else "file",
            "created_time": created,
            "modified_time": modified,
        }

        if path.is_file():
            node["size"] = size
            node["size_human"] = human_size(size)
            return node

        if depth >= max_depth:
            node["truncated"] = True
            node["children"] = []
            return node

        children = []
        try:
            for item in sorted(path.iterdir(), key=lambda x: (x.is_file(), x.name)):
                child_rel = f"{rel_clean}/{item.name}" if rel_clean else item.name
                children.append(build_tree(item, child_rel, depth + 1))
        except PermissionError:
            pass

        node["children"] = children
        return node

    return JSONResponse({"success": True, "tree": build_tree(target, rel_path)})


async def web_file_endpoint(request: Request):
    """Web UI endpoint for file preview."""
    web_cfg = get_web_config()
    auth_error = verify_web_auth(request, web_cfg)
    if auth_error:
        return auth_error

    file_path = request.query_params.get("path", "")
    if not file_path:
        return JSONResponse({"success": False, "error": "Missing path"}, status_code=400)

    user_data_dir = get_user_data_dir().resolve()
    target = (user_data_dir / file_path.lstrip("/")).resolve()

    if not str(target).startswith(str(user_data_dir)):
        return JSONResponse({"success": False, "error": "Access denied"}, status_code=403)
    if not target.exists():
        return JSONResponse({"success": False, "error": "File not found"}, status_code=404)
    if target.is_dir():
        return JSONResponse({"success": False, "error": "Path is a directory"}, status_code=400)

    max_preview_bytes = int(web_cfg.get("max_preview_bytes") or 200_000)
    try:
        file_size = target.stat().st_size
        truncated = file_size > max_preview_bytes
        read_bytes = max_preview_bytes if truncated else file_size
        with open(target, "rb") as f:
            content_bytes = f.read(read_bytes)
        content = content_bytes.decode("utf-8", errors="replace")
        return JSONResponse({
            "success": True,
            "content": content,
            "extension": target.suffix.lower(),
            "truncated": truncated,
            "name": target.name,
        })
    except Exception as e:
        return JSONResponse({"success": False, "error": str(e)}, status_code=500)


async def web_index(request: Request):
    """Serve web UI index.html."""
    web_cfg = get_web_config()
    if not web_cfg.get("enabled", False):
        return web_disabled_response()
    web_dir = Path(__file__).parent / "web"
    return FileResponse(web_dir / "index.html")


async def web_static(request: Request):
    """Serve web UI static files."""
    web_cfg = get_web_config()
    if not web_cfg.get("enabled", False):
        return web_disabled_response()
    rel = request.path_params.get("path") or ""
    web_dir = Path(__file__).parent / "web"
    target = (web_dir / rel).resolve()
    web_root = web_dir.resolve()
    if not str(target).startswith(str(web_root)) or not target.exists() or target.is_dir():
        target = web_root / "index.html"
    return FileResponse(target)


# ========== Route Registration ==========

def add_http_routes(app: Starlette) -> None:
    """Add all HTTP routes to the FastMCP app.

    This function should be called from the FilesystemFastMCP.http_app() method.
    """
    routes = [
        # Health/ready probes
        Route("/health", endpoint=healthcheck, methods=["GET"]),

        # Portal Homepage
        Route("/", endpoint=portal_homepage, methods=["GET"]),
        Route("/portal", endpoint=portal_homepage, methods=["GET"]),

        # Workspace Viewer
        Route("/workspace/viewer", endpoint=workspace_viewer, methods=["GET"]),

        # Admin API routes (must be before wildcard routes)
        Route("/admin/stats", endpoint=admin_stats, methods=["GET"]),
        Route("/admin/workspaces", endpoint=admin_list_workspaces, methods=["GET"]),
        Route("/admin/workspace/{workspace_id}/tree", endpoint=admin_workspace_tree, methods=["GET"]),
        Route("/admin/workspace/{workspace_id}", endpoint=admin_workspace_info, methods=["GET"]),
        Route("/admin/workspace/{workspace_id}", endpoint=admin_delete_workspace, methods=["DELETE"]),
        Route("/admin/cache/clear", endpoint=admin_clear_cache, methods=["POST"]),
        Route("/admin/previews", endpoint=admin_preview_stats, methods=["GET"]),

        # Admin Web UI (specific routes before wildcard)
        Route("/admin/api/tree", endpoint=web_tree_endpoint, methods=["GET"]),
        Route("/admin/api/file", endpoint=web_file_endpoint, methods=["GET"]),
        Route("/admin", endpoint=web_index, methods=["GET"]),
        Route("/admin/{path:path}", endpoint=web_static, methods=["GET"]),

        # User API routes
        Route("/api/workspace/tree", endpoint=user_workspace_tree, methods=["GET"]),
        Route("/api/workspace/file/download", endpoint=api_workspace_file_download, methods=["GET"]),
        Route("/api/workspace/file/preview", endpoint=api_workspace_file_preview, methods=["GET"]),
        Route("/api/workspace/file/upload", endpoint=api_workspace_file_upload, methods=["POST"]),
        Route("/api/workspace/file/update", endpoint=api_workspace_file_update, methods=["POST"]),
        Route("/api/workspace/file/delete", endpoint=api_workspace_file_delete, methods=["POST"]),
        Route("/api/workspace/download", endpoint=api_workspace_download, methods=["GET"]),
        Route("/api/workspace/exec", endpoint=api_exec_python, methods=["POST"]),
        # Image serving routes
        Route("/api/workspace/image/{path:path}", endpoint=api_workspace_image, methods=["GET"]),
        Route("/api/public/image/{filename:path}", endpoint=api_public_image, methods=["GET"]),

        # Preview Deployment API
        Route("/api/preview/deploy", endpoint=api_preview_deploy, methods=["GET", "POST"]),
        Route("/api/preview/status", endpoint=api_preview_status, methods=["GET"]),
        Route("/api/preview/stop", endpoint=api_preview_stop, methods=["GET", "POST"]),
        Route("/api/preview/access", endpoint=api_preview_access, methods=["GET", "POST"]),
        Route("/api/preview/mapping", endpoint=api_preview_mapping, methods=["GET"]),
    ]

    app.routes.extend(routes)
    logger.info(f"Added {len(routes)} HTTP routes to the application")
